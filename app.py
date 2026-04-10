import os
import re
import uuid
import json
import tempfile
import math
import requests
from datetime import datetime
from io import BytesIO
from flask import Flask, render_template, request, jsonify, Response, stream_with_context, send_file
from openai import OpenAI
from assistant_profile import ASSISTANT_PROFILE

try:
    from dotenv import load_dotenv
except Exception:
    load_dotenv = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from PIL import Image
except Exception:
    Image = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    RGBColor = None
    PP_ALIGN = None
    MSO_AUTO_SIZE = None
    MSO_VERTICAL_ANCHOR = None
    MSO_SHAPE = None

if load_dotenv:
    load_dotenv()

app = Flask(__name__)

# =========================
# PPT module imports
# Use the ppt/ package as the single source of truth for PPT generation.
# =========================
try:
    from ppt.generator import generate_ppt_structure as modular_generate_ppt_structure
except Exception:
    modular_generate_ppt_structure = None

try:
    from ppt.normalizer import normalize_ppt_json as modular_normalize_ppt_json
except Exception:
    modular_normalize_ppt_json = None

try:
    from ppt.normalizer import enforce_slide_content_budget as modular_enforce_slide_content_budget
except Exception:
    modular_enforce_slide_content_budget = None

try:
    from ppt.renderers import build_pptx_file as modular_build_pptx_file
    PPT_BUILDER_SOURCE = "ppt.renderers build_pptx_file"
except Exception:
    modular_build_pptx_file = None
    PPT_BUILDER_SOURCE = "app.py build_pptx_file"

try:
    from ppt.renderers import summarize_ppt_for_chat as modular_summarize_ppt_for_chat
except Exception:
    modular_summarize_ppt_for_chat = None


api_key = os.getenv("OPENROUTER_API_KEY")
if not api_key:
    raise ValueError("API key not found. Please set OPENROUTER_API_KEY in your environment variables.")

client = OpenAI(base_url="https://openrouter.ai/api/v1", api_key=api_key)
print("[startup] OPENROUTER_API_KEY loaded:", bool(api_key))
print("[startup] PPT_IMAGE_MODEL:", repr(os.getenv("PPT_IMAGE_MODEL")))
print("[startup] preferred ppt builder:", PPT_BUILDER_SOURCE)

conversation_histories = {}
uploaded_files = {}
generated_ppts = {}
conversation_resources = {}

DATA_DIR = os.path.join(tempfile.gettempdir(), "selena_app_data")
os.makedirs(DATA_DIR, exist_ok=True)
PPT_INDEX_PATH = os.path.join(DATA_DIR, "generated_ppts_index.json")
UPLOAD_INDEX_PATH = os.path.join(DATA_DIR, "uploaded_files_index.json")


def _load_json_file(path: str, default):
    try:
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception as e:
        print(f"[storage] failed to load {path}:", e)
    return default


def _save_json_file(path: str, payload):
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"[storage] failed to save {path}:", e)


def _persist_generated_ppts():
    _save_json_file(PPT_INDEX_PATH, generated_ppts)


def _persist_uploaded_files():
    _save_json_file(UPLOAD_INDEX_PATH, uploaded_files)


def _conversation_bucket(conversation_id: str) -> dict:
    return conversation_resources.setdefault(conversation_id, {"uploads": [], "ppts": []})


def _attach_upload_to_conversation(conversation_id: str, upload_id: str):
    if not upload_id:
        return
    bucket = _conversation_bucket(conversation_id)
    if upload_id not in bucket["uploads"]:
        bucket["uploads"].append(upload_id)


def _attach_ppt_to_conversation(conversation_id: str, file_id: str):
    if not file_id:
        return
    bucket = _conversation_bucket(conversation_id)
    if file_id not in bucket["ppts"]:
        bucket["ppts"].append(file_id)


def delete_generated_ppt_file(file_id: str) -> bool:
    info = generated_ppts.pop(file_id, None)
    if not info:
        return False
    try:
        path = info.get("path") or ""
        if path and os.path.exists(path):
            os.remove(path)
    except Exception as e:
        print(f"[ppt] failed to remove file {file_id}:", e)
    _persist_generated_ppts()
    return True


def delete_uploaded_file_record(upload_id: str) -> bool:
    removed = uploaded_files.pop(upload_id, None)
    if removed is not None:
        _persist_uploaded_files()
        return True
    return False


def clear_conversation_resources(conversation_id: str):
    bucket = conversation_resources.pop(conversation_id, {"uploads": [], "ppts": []})
    for upload_id in bucket.get("uploads", []):
        delete_uploaded_file_record(upload_id)
    for file_id in bucket.get("ppts", []):
        delete_generated_ppt_file(file_id)


generated_ppts.update(_load_json_file(PPT_INDEX_PATH, {}))
uploaded_files.update(_load_json_file(UPLOAD_INDEX_PATH, {}))


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SAFE_MARGIN_IN = 0.6
FOOTER_Y_IN = 6.82


def fit_single_text_to_box(text: str, box_width_in: float, box_height_in: float, preferred_size: int, min_size: int, max_len: int = None, bold: bool = False):
    txt = (text or '').strip()
    if max_len:
        txt = shrink_text(txt, max_len)
    if not txt:
        return '', preferred_size
    for size in range(preferred_size, min_size - 1, -1):
        needed_h = estimate_text_block_height([txt], box_width_in, size, line_spacing=1.15 if bold else 1.18, paragraph_gap_pt=0)
        if needed_h <= box_height_in:
            return txt, size
    if max_len:
        for cap in range(max(min(max_len, len(txt)), min_size*2), max(10, min(max_len, len(txt))//2), -4):
            shorter = shrink_text(txt, cap)
            needed_h = estimate_text_block_height([shorter], box_width_in, min_size, line_spacing=1.15 if bold else 1.18, paragraph_gap_pt=0)
            if needed_h <= box_height_in:
                return shorter, min_size
    return shrink_text(txt, max_len or max(18, int(box_width_in*10))), min_size


def populate_bullets(tf, bullet_points, size, color, space_after=6):
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    idx_out = 0
    for bullet in bullet_points or []:
        b = (bullet or "").strip()
        if not b:
            continue
        p = tf.paragraphs[0] if idx_out == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        style_paragraph(p, size=size, color=color)
        p.space_after = Pt(space_after)
        idx_out += 1


def safe_add_textbox(slide, left_in, top_in, width_in, height_in, text='', auto_fit=True):
    left_in = max(0.08, left_in)
    top_in = max(0.08, top_in)
    width_in = min(width_in, SLIDE_W_IN - left_in - 0.1)
    height_in = min(height_in, SLIDE_H_IN - top_in - 0.08)
    return add_textbox(slide, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in), text=text, auto_fit=auto_fit)

ALLOWED_EXTENSIONS = {
    ".py", ".txt", ".md", ".json", ".js", ".ts", ".html", ".css",
    ".java", ".cpp", ".c", ".cs", ".go", ".rs", ".php", ".xml",
    ".yaml", ".yml", ".pdf", ".docx"
}
MAX_FILE_CHARS = 18000


# =========================
# General helpers
# =========================
def get_history(conversation_id: str):
    if conversation_id not in conversation_histories:
        conversation_histories[conversation_id] = []
    return conversation_histories[conversation_id]


def get_mode_config(mode: str):
    modes = ASSISTANT_PROFILE.get("modes", {})
    return modes.get(mode, modes.get("chat", {}))


def build_system_prompt(mode: str) -> str:
    base_prompt = ASSISTANT_PROFILE["system_prompt"].strip()
    mode_prompt = get_mode_config(mode).get("prompt", "").strip()
    return base_prompt if not mode_prompt else base_prompt + "\n\n" + mode_prompt


def safe_filename(text: str, fallback: str = "selena_export") -> str:
    cleaned = re.sub(r'[\\/*?:"<>|]+', "_", (text or "").strip())
    cleaned = re.sub(r"\s+", "_", cleaned).strip("._")
    return cleaned[:60] or fallback


def generate_title_from_message(user_message: str) -> str:
    fallback = (user_message or "新聊天").strip()[:18] or "新聊天"
    try:
        response = client.chat.completions.create(
            model="openai/gpt-4o-mini",
            temperature=0.3,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You generate short conversation titles.\n"
                        "Rules:\n"
                        "- Return only the title text\n"
                        "- No quotes\n"
                        "- No punctuation unless necessary\n"
                        "- Keep it concise, ideally within 4 to 10 Chinese characters or 2 to 6 English words\n"
                        "- Match the user's language\n"
                    )
                },
                {
                    "role": "user",
                    "content": f"Generate a short title for this conversation:\n{user_message}"
                }
            ]
        )
        title = (response.choices[0].message.content or "").strip()
        return title[:30] if title else fallback
    except Exception as e:
        print("Title generation error:", e)
        return fallback


def infer_mode_by_filename(filename: str) -> str:
    lower = filename.lower()
    code_suffixes = (
        ".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go",
        ".rs", ".php", ".html", ".css", ".xml", ".json", ".yaml", ".yml"
    )
    if lower.endswith(code_suffixes):
        return "code"
    return "chat"


# =========================
# File parsing
# =========================
def extract_text_from_pdf(file_bytes: bytes) -> str:
    if PdfReader is None:
        raise RuntimeError("pypdf is not installed.")
    reader = PdfReader(BytesIO(file_bytes))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    return "\n\n".join(parts).strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    if Document is None:
        raise RuntimeError("python-docx is not installed.")
    doc = Document(BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()]).strip()


def read_uploaded_file(filename: str, file_bytes: bytes) -> str:
    extension = os.path.splitext(filename)[1].lower()
    text_exts = {
        ".py", ".txt", ".md", ".json", ".js", ".ts", ".html", ".css",
        ".java", ".cpp", ".c", ".cs", ".go", ".rs", ".php", ".xml",
        ".yaml", ".yml"
    }
    if extension in text_exts:
        return file_bytes.decode("utf-8", errors="replace")
    if extension == ".pdf":
        return extract_text_from_pdf(file_bytes)
    if extension == ".docx":
        return extract_text_from_docx(file_bytes)
    raise ValueError("Unsupported file type.")


def build_file_augmented_prompt(file_obj: dict, base_message: str, mode: str) -> str:
    filename = file_obj["filename"]
    extension = file_obj["extension"].lstrip(".") or "txt"
    content = file_obj["content"]

    if mode == "code":
        request_text = base_message or "请分析这个文件的结构、功能、潜在问题，并给出改进建议。"
        return (
            f"我上传了一个文件，请你以代码助手的方式帮我处理。\n\n"
            f"文件名：{filename}\n"
            f"类型：{extension}\n\n"
            f"我的要求：{request_text}\n\n"
            f"以下是文件内容：\n\n{content}"
        )

    if mode == "ppt":
        request_text = base_message or "请基于这份内容为我整理一个清晰的 PPT 大纲和每页要点。"
        return (
            f"我上传了一个资料文件，请你以 PPT 助手的方式帮我处理。\n\n"
            f"文件名：{filename}\n"
            f"类型：{extension}\n\n"
            f"我的要求：{request_text}\n\n"
            f"以下是文件内容：\n\n{content}"
        )

    request_text = base_message or "请总结这份文件，并告诉我重点内容。"
    return (
        f"我上传了一个文件，请先阅读并理解它。\n\n"
        f"文件名：{filename}\n"
        f"类型：{extension}\n\n"
        f"我的要求：{request_text}\n\n"
        f"以下是文件内容：\n\n{content}"
    )



def _walk_json(node):
    if isinstance(node, dict):
        yield node
        for value in node.values():
            yield from _walk_json(value)
    elif isinstance(node, list):
        for item in node:
            yield from _walk_json(item)


def extract_image_urls(result_json: dict) -> list[str]:
    image_urls = []
    markdown_image_pattern = re.compile(r'!\[[^\]]*\]\((https?://[^)]+|data:image/[^)]+)\)')
    plain_url_pattern = re.compile(r'(https?://\S+\.(?:png|jpg|jpeg|webp)(?:\?\S*)?|data:image/[^\s"\']+)')

    for node in _walk_json(result_json):
        if isinstance(node, dict):
            direct = node.get("url")
            if isinstance(direct, str) and direct.startswith(("http://", "https://", "data:image/")):
                image_urls.append(direct)

            for key in ("image_url", "imageUrl"):
                nested = node.get(key)
                if isinstance(nested, dict):
                    url = nested.get("url")
                    if isinstance(url, str) and url.startswith(("http://", "https://", "data:image/")):
                        image_urls.append(url)

        elif isinstance(node, str):
            for match in markdown_image_pattern.findall(node):
                image_urls.append(match.strip())
            for match in plain_url_pattern.findall(node):
                image_urls.append(match.strip().rstrip('.,)'))

    seen = set()
    unique_urls = []
    for url in image_urls:
        if url not in seen:
            seen.add(url)
            unique_urls.append(url)
    return unique_urls


def extract_image_b64_strings(result_json: dict) -> list[str]:
    candidates = []
    b64_keys = {
        "b64_json", "b64", "base64", "image_base64", "imageBase64",
        "png", "jpeg", "jpg", "webp"
    }
    for node in _walk_json(result_json):
        if not isinstance(node, dict):
            continue
        for key, value in node.items():
            if key in b64_keys and isinstance(value, str) and len(value) > 200:
                candidates.append(value.strip())
    deduped = []
    seen = set()
    for item in candidates:
        if item not in seen:
            seen.add(item)
            deduped.append(item)
    return deduped


def save_base64_image(b64_data: str, suffix: str = ".png") -> str:
    import base64
    raw = b64_data.strip()
    if raw.startswith("data:image/"):
        header, raw = raw.split(",", 1)
        if "jpeg" in header or "jpg" in header:
            suffix = ".jpg"
        elif "webp" in header:
            suffix = ".webp"
        else:
            suffix = ".png"
    raw += "=" * (-len(raw) % 4)
    binary = base64.b64decode(raw)
    fd, path = tempfile.mkstemp(prefix="selena_img_", suffix=suffix)
    os.close(fd)
    with open(path, "wb") as f:
        f.write(binary)
    return path


def download_or_decode_image(source: str, suffix: str = ".png") -> str:
    source = (source or "").strip()
    if not source:
        return ""
    if source.startswith("data:image/"):
        return save_base64_image(source, suffix)
    if source.startswith(("http://", "https://")):
        return download_binary_file(source, suffix)
    raise ValueError("Unsupported image source format.")


# =========================
# PPT generation
# =========================
def normalize_ppt_json(data: dict) -> dict:
    theme = str(data.get("theme") or "selena").lower()
    if theme not in {"selena", "business", "academic"}:
        theme = "selena"
    slides = data.get("slides") or []
    normalized_slides = []
    for item in slides:
        if not isinstance(item, dict):
            continue
        slide_type = str(item.get("type") or "bullet").lower()
        if slide_type not in {"section", "bullet", "highlight", "two_column", "compare", "image_left", "image_right", "timeline"}:
            slide_type = "bullet"
        bullets = item.get("bullets") or []
        if not isinstance(bullets, list):
            bullets = []
        timeline_steps = item.get("timeline_steps") or []
        if not isinstance(timeline_steps, list):
            timeline_steps = []
        normalized_slides.append({
            "type": slide_type,
            "title": str(item.get("title") or "").strip(),
            "summary": str(item.get("summary") or "").strip(),
            "bullets": [str(x).strip() for x in bullets if str(x).strip()],
            "timeline_steps": [str(x).strip() for x in timeline_steps if str(x).strip()],
            "left_title": str(item.get("left_title") or "").strip(),
            "right_title": str(item.get("right_title") or "").strip(),
            "left_points": [str(x).strip() for x in (item.get("left_points") or []) if str(x).strip()],
            "right_points": [str(x).strip() for x in (item.get("right_points") or []) if str(x).strip()],
            "highlight": str(item.get("highlight") or "").strip(),
            "closing": str(item.get("closing") or "").strip(),
            "speaker_note": str(item.get("speaker_note") or "").strip(),
            "image_prompt": str(item.get("image_prompt") or "").strip(),
            "image_caption": str(item.get("image_caption") or "").strip(),
            "image_path": "",
        })
    data["theme"] = theme
    data["slides"] = normalized_slides
    return data


def generate_ppt_structure(topic: str, upload_id: str = "") -> dict:
    uploaded = uploaded_files.get(upload_id) if upload_id else None
    file_context = ""
    if uploaded:
        file_context = (
            f"\n\n参考资料文件名：{uploaded['filename']}\n"
            f"参考资料内容：\n{uploaded['content'][:12000]}"
        )

    prompt = f"""请根据下面的主题生成一份“可直接上台演讲、内容扎实、并包含视觉页”的 PPT 结构，输出严格 JSON。
主题：{topic}
{file_context}

你需要自动判断 theme，可选值只有：selena / business / academic。
你需要为每一页自动选择 type，可选值：section / bullet / highlight / two_column / compare / image_left / image_right / timeline。
- timeline：用于路线图、阶段计划、里程碑、排期；使用 timeline_steps（3~5 条，每条为完整句，含时间或阶段信息）；可同时保留 bullets 与 timeline_steps 一致以便兼容。
第一页封面不用放进 slides，系统会自动生成封面。

【核心目标】
1. 这不是“泛泛而谈”的大纲，而是可以直接讲的演示稿内容。
2. 每一页都要有一个明确结论（takeaway），summary 必须体现判断，不要只做背景介绍。
3. 全文保持“结论 -> 依据 -> 行动”的表达习惯，确保听众能听完就知道下一步做什么。
4. 对于适合视觉化表达的内容，必须主动使用 image_left 或 image_right。
5. 若用户需求较短，可合理补全结构，但必须避免空泛套话，优先给出可执行、可落地的信息。

【视觉页要求】
- 整份 PPT 至少包含 1 页、最多 2 页 image_left 或 image_right。
- 每个视觉页都必须给出 image_prompt，写成适合生成演示图片的英文提示词。
- image_prompt 要具体，包含主体、场景、风格、配色或构图，避免抽象空词。
- image_caption 用一句简短中文说明图片如何支撑本页结论，不能只写“示意图”。
- 图片应优先适配 16:9 商务演示风格，避免复杂背景、过多人物、文字水印。

【内容质量与反空话要求】
- summary：1~2 句，必须包含明确判断或结论，中文建议 28~60 字。
- bullet / left_points / right_points：每条中文建议 18~38 字，必须是完整表达，并至少包含以下要素之一：动作 / 依据 / 指标 / 时间 / 风险。
- 禁止高频空话单独成句，如“全面提升效率”“持续赋能业务”“增强竞争力”等；如果出现，必须配具体对象和落地方式。
- speaker_note：2~4 句，至少包含这三类信息中的两类：如何讲这一页、要强调的数据或证据、与下一页的过渡。
- highlight：必须是可直接上台强调的核心判断句，避免口号化。
- 如果缺少可验证数据，不要编造数字；可使用“假设：”或“待补充：”明确标注。

【去重与完整性】
- 每页 summary 不得与 title 仅作同义复述或简单扩写；须额外给出判断、范围或收束信息。
- 同页 bullet / left_points / right_points 不得与 title、summary 语义重复（勿换说法复述同一句）。
- image_left / image_right：除页顶 summary 外，必须提供至少 2 条 bullets，写行动、依据或下一步，不得复述页顶 summary 原句。
- image_caption 专门说明“图示如何支撑本页结论”，禁止照抄 summary。
- 禁止在 title、summary、bullet、caption、highlight、closing 等可见字段末尾使用「…」或「...」；须用完整句号收束。
- 顶层 subtitle（封面副标题）：须 **40~95 字**，建议写成 2~3 句（可用分号衔接）：①汇报对象与场景；②材料覆盖范围或时间；③本报告要回答的一个核心问题。禁止只做主标题的同义压缩。
- section 页（如引言与背景）：summary 须 **85~220 字**，**至少四句**，依次写清：行业与组织背景、关键趋势或数据感知、当前主要矛盾或机会、本节将交付的结论与后文结构（勿写成标题扩写）。
- bullet 页：bullets 至少 3 条、优先 4 条；不得输出空字符串占位。
- compare 页：left_points 与 right_points 各至少 3 条、至多 4 条，形成一一对照。
- highlight 页：closing 必须与 highlight 形成递进（行动号召、下一步或风险提醒），禁止与 highlight 同义改写。
- 视觉页：bullets 为 2~4 条非空完整句；image_caption 从“图里有什么、如何印证结论”角度写，避免复述页顶 summary 的核心谓语。
- 结论类视觉页：image_caption 与 bullets **禁止**元话术（如「这一页总结了…」「本页旨在…」「报告的核心结论」等）；须写**具体行动或证据**（负责人、时间窗口、交付物、指标之一）。
- image_caption 禁止元话术（如「总结报告的核心要点」「本页将展示」等）和半截句；须写具体结论关系，句末只能是句号问号叹号。
- image_prompt（英文）须强调画面无任何文字与标题条，仅用隐喻或场景表达，避免出现 slide title / headline 构图。

【结构要求】
- 总页数控制在 7 到 9 页（不含封面）。
- 至少包含 1 页 highlight。
- 至少包含 1 页 two_column 或 compare。
- 至少包含 2 页 bullet。
- 若叙事含落地计划、路线图、阶段目标，至少 1 页使用 timeline（勿用普通 bullet 凑合）。
- 建议整体叙事遵循：背景与问题 -> 关键洞察 -> 方案与路径 -> 落地计划 -> 结论与行动。
- 不要生成空字段，不要生成只有一两个字的内容，不要重复同义内容。

JSON 格式必须是：
{{
"title":"总标题",
"subtitle":"副标题",
"theme":"selena",
"slides":[
    {{
    "type":"section / bullet / highlight / two_column / compare / image_left / image_right / timeline",
    "title":"页标题",
    "summary":"这一页的核心结论句或两句简要说明，不能为空",
    "bullets":["完整表达1","完整表达2","完整表达3"],
    "left_title":"左侧标题，可选",
    "right_title":"右侧标题，可选",
    "left_points":["左侧完整表达1","左侧完整表达2","左侧完整表达3"],
    "right_points":["右侧完整表达1","右侧完整表达2","右侧完整表达3"],
    "highlight":"适用于 highlight 页的大句子，必须有判断性",
    "closing":"适用于 highlight 页的收束句",
    "speaker_note":"这一页该怎么讲，2到4句",
    "image_prompt":"用于生成页面图片的英文提示词，视觉页必填",
    "image_caption":"图片说明，中文简短一句",
    "timeline_steps":["阶段或里程碑1","阶段或里程碑2","阶段或里程碑3"]
    }}
]
}}

输出规则：
1. 只返回 JSON。
2. 不要 markdown，不要代码块，不要额外解释。
3. 商务主题偏结论、价值和机会；学术主题偏逻辑、方法和结论；产品主题偏流程、架构和体验。
4. 视觉页中的图片要与页面核心观点紧密相关，而不是泛泛而谈。
"""

    response = client.chat.completions.create(
        model="openai/gpt-4o-mini",
        temperature=0.6,
        messages=[
            {
                "role": "system",
                "content": (
                    "You create high-quality, speech-ready PPT structures and return strict JSON only. "
                    "The slides must be specific, actionable, and non-generic, with clear takeaway-driven logic. "
                    "Include 1-2 visual slides with concrete English image prompts that directly support slide conclusions."
                )
            },
            {"role": "user", "content": prompt}
        ]
    )

    raw = (response.choices[0].message.content or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:].strip()

    data = json.loads(raw)
    if "slides" not in data or not isinstance(data["slides"], list) or not data["slides"]:
        raise ValueError("Invalid PPT JSON structure.")
    data = normalize_ppt_json(data)
    data = refine_ppt_structure_quality(topic=topic, ppt_data=data, has_upload=bool(uploaded))
    return data


def refine_ppt_structure_quality(topic: str, ppt_data: dict, has_upload: bool = False) -> dict:
    slides = ppt_data.get("slides") or []
    if not slides:
        return ppt_data

    review_prompt = f"""你将收到一份已生成的 PPT JSON。你的任务是“质量自检并仅重写薄弱页面”，然后返回完整 JSON。

主题：{topic}
是否有参考资料：{"是" if has_upload else "否"}

【目标】
把内容从“能看”提升为“可直接上台讲”，重点修正空泛、重复、口号化表达。

【自检标准】
1. 每页都必须有明确结论（summary 体现判断，而非背景介绍）。
2. bullet / left_points / right_points 不能只是正确的废话；每条尽量包含动作、依据、指标、时间、风险中的至少一种信息。
3. speaker_note 要有讲法：至少体现“怎么讲/强调什么/如何过渡”中的两项。
4. 不能编造具体数据；若缺数据可写“假设：”或“待补充：”。
5. 视觉页 image_caption 必须解释图片与本页结论的关系；image_prompt 必须具体且可生成，不要抽象空词堆砌。
6. 保持原有页面数量、页面顺序、type 分布（尤其保留 1~2 页视觉页；若有 timeline 页须保留 timeline 类型与 timeline_steps），只优化内容质量。
7. 消除 title 与 summary 的重复；视觉页右侧正文禁止再复述页顶 summary（可改为行动项、证据或下一步）。
8. image_left / image_right 至少保留 2 条不重复的 bullets，且不得与 image_caption、summary 同义重复。
9. 去掉可见字段末尾的「…」「...」，改为完整句末标点。
10. 顶层 subtitle 若少于约 40 字或与 title 同义，须扩展为 40~95 字，包含受众、范围与核心问题。
11. section 页 summary 若少于约 85 字或不足四句，须扩展为 85~220 字（背景—趋势—矛盾/机会—本节交付与结构）。
12. bullet 页若不足 3 条，须补全；compare 每侧若不足 3 条，须补全；不得留下空字符串 bullet。
13. highlight 的 closing 若与 highlight 同义，须改为行动项、下一步或风险提醒。
14. 视觉页 bullets 中不得出现空串；caption 与页顶 summary 谓语高度重叠时须改写 caption 角度。
15. 视觉页 image_caption 不得与任一条 bullet 开头整句重复；不得出现以「...」或「…」结尾的未写完句。
16. 若 image_caption 含元话术或半成品，须改为具体行动或证据句；image_prompt 须禁止图中出现任何文字与英文标题。
17. 若 image_caption 或 bullet 以「这一页/本页…总结」开头，须整句改写为具体结论或行动项，不得保留幻灯片说明语。
18. 禁止出现多页 timeline（或 bullet 升格为时间轴后）**标题、summary、timeline_steps 完全一致**的重复页；若发现重复，只保留一页并将其余改为不同议题的 bullet/compare 或合并为单页更完整的时间轴。

【重写策略】
- 仅重写低质量字段：summary / bullets / left_points / right_points / highlight / closing / speaker_note / image_caption / image_prompt。
- 对高质量字段尽量保持不变，避免整份重写造成风格漂移。

【输入 JSON】
{json.dumps(ppt_data, ensure_ascii=False)}

【输出规则】
1. 只返回 JSON。
2. 不要 markdown，不要解释，不要代码块。
3. 返回结构必须与输入 schema 一致，且 slides 不得为空。
"""

    try:
        review_response = client.chat.completions.create(
            model="openai/gpt-4o-mini",
            temperature=0.3,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a strict PPT quality reviewer and rewriter. "
                        "Improve weak slide content while preserving structure and order. "
                        "Return strict JSON only."
                    ),
                },
                {"role": "user", "content": review_prompt},
            ],
        )
        reviewed_raw = (review_response.choices[0].message.content or "").strip()
        if reviewed_raw.startswith("```"):
            reviewed_raw = reviewed_raw.strip("`")
            if reviewed_raw.lower().startswith("json"):
                reviewed_raw = reviewed_raw[4:].strip()
        reviewed_data = json.loads(reviewed_raw)
        if "slides" in reviewed_data and isinstance(reviewed_data["slides"], list) and reviewed_data["slides"]:
            return normalize_ppt_json(reviewed_data)
    except Exception as e:
        print("PPT quality refine fallback:", e)
    return ppt_data


def build_visual_prompt(topic: str, slide_data: dict, theme: str) -> str:
    theme = (theme or "selena").lower()
    style_map = {
        "business": "premium business presentation illustration, corporate photography style, clean lighting, navy and blue palette, minimal background, 16:9 slide composition, no text, no watermark",
        "academic": "clean academic presentation illustration, professional infographic style, white and blue palette, clear structure, 16:9 slide composition, no text, no watermark",
        "selena": "futuristic presentation illustration, modern product design style, blue purple gradient palette, minimal cinematic lighting, 16:9 slide composition, no text, no watermark",
    }
    bullets = slide_data.get("bullets") or []
    detail = "; ".join(bullets[:3])
    summary = slide_data.get("summary") or ""
    title = slide_data.get("title") or topic or "presentation"
    base = f"{title}. {summary}. {detail}".strip()
    base = re.sub(r"\s+", " ", base)
    return f"{base}. {style_map.get(theme, style_map['selena'])}"


def llm_generate_image_prompt(topic: str, slide_data: dict, theme: str) -> str:
    try:
        prompt = (
            "Generate one concise English image prompt for a PowerPoint slide. Return only the prompt text. "
            "The image must be presentation-friendly, visually clean, with no text of any kind (no letters, captions, titles, labels), "
            "no watermark, no slide-title graphics, metaphor or scene only, suitable for a 16:9 slide.\n\n"
            f"Theme: {theme}\n"
            f"Presentation topic: {topic}\n"
            f"Slide title: {slide_data.get('title', '')}\n"
            f"Slide summary: {slide_data.get('summary', '')}\n"
            f"Slide bullets: {' | '.join(slide_data.get('bullets') or [])}\n"
            f"Image caption hint: {slide_data.get('image_caption', '')}\n"
        )
        response = client.chat.completions.create(
            model="openai/gpt-4o-mini",
            temperature=0.4,
            messages=[
                {"role": "system", "content": "You write strong English prompts for slide images."},
                {"role": "user", "content": prompt},
            ],
        )
        text = (response.choices[0].message.content or "").strip()
        text = text.strip('`').strip()
        if text:
            return text
    except Exception as e:
        print("Image prompt generation fallback:", e)
    return build_visual_prompt(topic, slide_data, theme)


def ensure_visual_slides(ppt_data: dict, topic: str, min_images: int = 1, max_images: int = 2) -> dict:
    slides = ppt_data.get("slides") or []
    visual_slides = [s for s in slides if s.get("type") in {"image_left", "image_right"}]

    # backfill missing prompts on existing visual slides
    for idx, slide in enumerate(visual_slides):
        if not slide.get("image_prompt"):
            slide["image_prompt"] = llm_generate_image_prompt(topic, slide, ppt_data.get("theme", "selena"))
        if not slide.get("image_caption"):
            slide["image_caption"] = slide.get("summary") or slide.get("title") or "页面视觉图"
        slide["type"] = "image_left" if idx % 2 == 0 else "image_right"

    if len(visual_slides) >= min_images:
        return ppt_data

    candidates = []
    for i, slide in enumerate(slides):
        score = choose_visual_candidate_score(slide)
        if score > -999:
            candidates.append((score, i, slide))
    candidates.sort(reverse=True)

    needed = min(min_images, max_images) - len(visual_slides)
    converted = 0
    for _, _, slide in candidates:
        if converted >= needed:
            break
        if slide.get("type") in {"image_left", "image_right"}:
            continue
        slide["type"] = "image_left" if converted % 2 == 0 else "image_right"
        if not slide.get("image_caption"):
            slide["image_caption"] = slide.get("summary") or slide.get("title") or "页面视觉图"
        slide["image_prompt"] = llm_generate_image_prompt(topic, slide, ppt_data.get("theme", "selena"))
        converted += 1

    return ppt_data


def choose_ppt_image_model() -> str:
    image_config = ASSISTANT_PROFILE.get("image_generation", {})
    return image_config.get("model", "google/gemini-2.5-flash-image")


def generate_presentation_image_assets(prompt: str, aspect_ratio: str = "16:9", image_size: str = "1K") -> dict:
    if not prompt:
        return {"urls": [], "b64": [], "raw": {}}

    payload = {
        "model": choose_ppt_image_model(),
        "messages": [{
            "role": "user",
            "content": (
                "Create a clean, presentation-ready image for a slide. "
                "Strictly no text in the image: no letters, words, titles, captions, labels, or logos in any language. "
                "No watermarks or UI. Use metaphor or scene only.\n\n"
                f"{prompt}"
            )
        }],
        "modalities": ["image", "text"],
        "stream": False,
        "image_config": {"aspect_ratio": aspect_ratio, "image_size": image_size}
    }

    response = requests.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json=payload,
        timeout=240
    )
    response.raise_for_status()
    result = response.json()
    urls = extract_image_urls(result)
    b64 = extract_image_b64_strings(result)
    print("Image API returned urls:", len(urls), "b64:", len(b64))
    if not urls and not b64:
        print("Image API raw keys:", list(result.keys())[:10])
        print("Image API raw preview:", json.dumps(result, ensure_ascii=False)[:1200])
    return {"urls": urls, "b64": b64, "raw": result}


def try_generate_slide_image_file(prompt: str) -> str:
    attempts = [("16:9", "1K"), ("4:3", "1K"), ("1:1", "1K")]
    for aspect_ratio, image_size in attempts:
        try:
            assets = generate_presentation_image_assets(prompt=prompt, aspect_ratio=aspect_ratio, image_size=image_size)
            path = ""
            if assets["urls"]:
                first = assets["urls"][0]
                print("Using image url/data source:", first[:120])
                path = download_or_decode_image(first, ".png")
            elif assets["b64"]:
                print("Using base64 image payload")
                path = save_base64_image(assets["b64"][0], ".png")
            if path and is_image_file_usable(path):
                return path
            if path:
                print("Generated image rejected by validator:", path)
        except Exception as e:
            print(f"Image generation attempt failed ({aspect_ratio}, {image_size}):", e)
    return ""


def download_binary_file(url: str, suffix: str = ".png") -> str:
    resp = requests.get(url, timeout=180)
    resp.raise_for_status()
    fd, path = tempfile.mkstemp(prefix="selena_img_", suffix=suffix)
    os.close(fd)
    with open(path, "wb") as f:
        f.write(resp.content)
    return path


def ensure_slide_images(ppt_data: dict, topic: str = "", max_images: int = 2) -> dict:
    ppt_data = ensure_visual_slides(
        ppt_data,
        topic=topic or ppt_data.get("title", "Presentation"),
        min_images=1,
        max_images=max_images
    )
    generated = 0
    for slide in ppt_data.get("slides", []):
        if generated >= max_images:
            break
        slide_type = slide.get("type") or "bullet"
        prompt = (slide.get("image_prompt") or "").strip()
        if slide_type not in {"image_left", "image_right"}:
            continue
        if not prompt:
            prompt = llm_generate_image_prompt(
                topic or ppt_data.get("title", "Presentation"),
                slide,
                ppt_data.get("theme", "selena")
            )
            slide["image_prompt"] = prompt
        try:
            image_path = try_generate_slide_image_file(prompt)
            print("Slide image path:", image_path or "[none]")
            if image_path:
                slide["image_path"] = image_path
                generated += 1
            else:
                slide["image_path"] = ""
        except Exception as e:
            print("Slide image generation error:", e)
            slide["image_path"] = ""
    ppt_data["generated_image_count"] = generated
    return ppt_data


def add_picture_fit(slide, image_path, left, top, width, height, padding=Inches(0.08), mode="contain"):
    if not image_path or not os.path.exists(image_path):
        return None

    inner_left = max(left + padding, Inches(0.05))
    inner_top = max(top + padding, Inches(0.05))
    max_right = Inches(SLIDE_W_IN - 0.08)
    max_bottom = Inches(SLIDE_H_IN - 0.08)
    inner_width = max(min(width - padding * 2, max_right - inner_left), Inches(0.2))
    inner_height = max(min(height - padding * 2, max_bottom - inner_top), Inches(0.2))

    if Image is None:
        return slide.shapes.add_picture(image_path, inner_left, inner_top, width=inner_width, height=inner_height)

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    if not img_w or not img_h:
        return None

    box_ratio = float(inner_width) / float(inner_height) if float(inner_height) else 1.0
    img_ratio = img_w / img_h
    mode = (mode or "contain").lower()

    if mode == "cover":
        if img_ratio > box_ratio:
            pic_height = inner_height
            pic_width = min(inner_height * img_ratio, Inches(SLIDE_W_IN - 0.16))
            pic_left = max(Inches(0.05), inner_left - (pic_width - inner_width) / 2)
            pic_top = inner_top
        else:
            pic_width = inner_width
            pic_height = min(inner_width / img_ratio, Inches(SLIDE_H_IN - 0.16))
            pic_left = inner_left
            pic_top = max(Inches(0.05), inner_top - (pic_height - inner_height) / 2)
    else:
        if img_ratio > box_ratio:
            pic_width = inner_width
            pic_height = inner_width / img_ratio if img_ratio else inner_height
            pic_left = inner_left
            pic_top = inner_top + (inner_height - pic_height) / 2
        else:
            pic_height = inner_height
            pic_width = inner_height * img_ratio
            pic_left = inner_left + (inner_width - pic_width) / 2
            pic_top = inner_top

    pic_width = min(pic_width, max_right - pic_left)
    pic_height = min(pic_height, max_bottom - pic_top)
    return slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width, height=pic_height)


def get_theme_tokens(theme: str) -> dict:
    theme = (theme or "selena").lower()
    if theme == "business":
        return {
            "primary": RGBColor(15, 23, 42),
            "secondary": RGBColor(37, 99, 235),
            "accent": RGBColor(219, 234, 254),
            "muted": RGBColor(71, 85, 105),
            "light": RGBColor(248, 250, 252),
            "line": RGBColor(191, 219, 254),
        }
    if theme == "academic":
        return {
            "primary": RGBColor(30, 64, 175),
            "secondary": RGBColor(14, 116, 144),
            "accent": RGBColor(224, 242, 254),
            "muted": RGBColor(71, 85, 105),
            "light": RGBColor(255, 255, 255),
            "line": RGBColor(186, 230, 253),
        }
    return {
        "primary": RGBColor(79, 70, 229),
        "secondary": RGBColor(139, 92, 246),
        "accent": RGBColor(237, 233, 254),
        "muted": RGBColor(71, 85, 105),
        "light": RGBColor(255, 255, 255),
        "line": RGBColor(196, 181, 253),
    }


def add_full_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", auto_fit=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_fit else MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    if text:
        tf.text = text
    return box, tf


def style_paragraph(paragraph, size=20, bold=False, color=None, align=None):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def clamp_text(text: str, max_len: int = 120) -> str:
    text = (text or "").strip()
    if len(text) <= max_len:
        return text
    candidate = text[:max_len]
    cut_chars = ["。", "！", "？", ".", "!", "?", "；", ";", "，", ",", ":", "：", "、"]
    last_pos = -1
    for ch in cut_chars:
        pos = candidate.rfind(ch)
        if pos > last_pos:
            last_pos = pos
    if last_pos >= max(8, int(max_len * 0.55)):
        return candidate[: last_pos + 1].strip()
    shortened = candidate.rstrip("，,。.!！?？:：;；、 ")
    if not shortened:
        shortened = candidate.strip()
    has_cjk = any("\u4e00" <= c <= "\u9fff" for c in shortened)
    if has_cjk:
        if shortened[-1] not in "。！？":
            return shortened + "。"
        return shortened
    if shortened[-1] not in ".!?":
        return shortened + "."
    return shortened


def _slide_body_dedupe_key(s: str) -> str:
    return " ".join((s or "").replace("•", "").replace("·", "").split()).strip().lower()


def adaptive_font_size(text: str, default: int, minimum: int, breakpoints=None) -> int:
    text = (text or "").strip()
    breakpoints = breakpoints or []
    size = default
    for char_count, candidate in sorted(breakpoints):
        if len(text) > char_count:
            size = min(size, candidate)
    return max(minimum, size)


def prepare_points(points, fallback_text="", max_items: int = 4, max_len: int = 54):
    result = []
    for item in points or []:
        item = clamp_text(str(item), max_len)
        if item:
            result.append(item)
        if len(result) >= max_items:
            break
    if not result and fallback_text:
        result.append(clamp_text(fallback_text, max_len))
    return result


def detect_text_language(text: str) -> str:
    text = text or ""
    zh_count = sum(1 for ch in text if '一' <= ch <= '鿿')
    return "zh" if zh_count >= max(2, len(text) * 0.25) else "en"


def estimate_lines(text: str, box_width_in: float, font_size_pt: int, prefix_chars: float = 0.0) -> int:
    text = (text or "").strip()
    if not text:
        return 1
    lang = detect_text_language(text)
    avg_char_factor = 0.95 if lang == "zh" else 0.56
    box_width_pt = max(24.0, box_width_in * 72.0)
    chars_per_line = max(1.0, box_width_pt / (font_size_pt * avg_char_factor))
    effective_len = len(text) + prefix_chars
    return max(1, int(math.ceil(effective_len / chars_per_line)))


def estimate_text_block_height(texts, box_width_in: float, font_size_pt: int, line_spacing=1.22, paragraph_gap_pt=6, prefix_chars: float = 0.0) -> float:
    total_pt = 0.0
    items = list(texts or [])
    if not items:
        return font_size_pt * line_spacing / 72.0
    for idx, text in enumerate(items):
        lines = estimate_lines(text, box_width_in, font_size_pt, prefix_chars=prefix_chars)
        total_pt += lines * font_size_pt * line_spacing
        if idx < len(items) - 1:
            total_pt += paragraph_gap_pt
    return total_pt / 72.0


def shrink_text(text: str, max_len: int) -> str:
    text = re.sub(r'\s+', ' ', (text or '').strip())
    text = re.sub(r'^(首先|其次|另外|此外|总的来说|总体来看|需要注意的是)[，,:： ]*', '', text)
    text = re.sub(r'(可以看到|我们可以看到|这说明|这意味着)[，,:： ]*', '', text)
    return clamp_text(text, max_len)


def fit_bullets_to_box(points, box_width_in, box_height_in, preferred_size=18, min_size=14,
                       max_items=4, max_len=54, fallback_text=""):
    cleaned = []
    for item in points or []:
        txt = str(item).strip()
        if txt:
            cleaned.append(txt)
    if not cleaned and fallback_text:
        cleaned = [fallback_text]

    cleaned = [shrink_text(x, max_len) for x in cleaned[:max_items]]
    cleaned = [x for x in cleaned if isinstance(x, str) and x.strip()]
    if not cleaned and fallback_text:
        fb = shrink_text(str(fallback_text).strip(), max_len)
        if fb and fb.strip():
            cleaned = [fb]
    if not cleaned:
        return [], preferred_size, False

    for size in range(preferred_size, min_size - 1, -1):
        needed_h = estimate_text_block_height(cleaned, box_width_in, size, line_spacing=1.18, paragraph_gap_pt=7, prefix_chars=2.2)
        if needed_h <= box_height_in:
            return cleaned, size, False

    trimmed = cleaned[:]
    reduced = False
    while len(trimmed) > 2:
        trimmed = trimmed[:-1]
        reduced = True
        for size in range(preferred_size, min_size - 1, -1):
            needed_h = estimate_text_block_height(trimmed, box_width_in, size, line_spacing=1.18, paragraph_gap_pt=7, prefix_chars=2.2)
            if needed_h <= box_height_in:
                return trimmed, size, reduced

    compact = [shrink_text(x, max(16, max_len - 10)) for x in trimmed]
    compact = [x for x in compact if isinstance(x, str) and x.strip()]
    for size in range(preferred_size, min_size - 1, -1):
        needed_h = estimate_text_block_height(compact, box_width_in, size, line_spacing=1.16, paragraph_gap_pt=5, prefix_chars=2.2)
        if needed_h <= box_height_in:
            return compact, size, True

    return compact[:2], min_size, True


def choose_visual_candidate_score(slide: dict) -> int:
    slide_type = slide.get("type") or "bullet"
    if slide_type not in {"bullet", "highlight", "two_column", "compare"}:
        return -999
    score = 0
    summary = slide.get("summary") or ""
    bullets = slide.get("bullets") or []
    total_points = len(bullets) + len(slide.get("left_points") or []) + len(slide.get("right_points") or [])
    score += 12 if slide_type == "highlight" else 0
    score += 6 if 18 <= len(summary) <= 58 else 0
    score += 5 if total_points <= 3 else 0
    score -= 5 if total_points >= 6 else 0
    score -= 8 if slide_type == "compare" else 0
    score += 4 if slide.get("highlight") else 0
    return score


def enforce_slide_content_budget(ppt_data: dict) -> dict:
    diagnostics = []
    new_slides = []
    for slide in ppt_data.get("slides", []):
        slide = dict(slide)
        st = slide.get("type") or "bullet"
        slide["title"] = clamp_text(slide.get("title") or "", 36)
        slide["summary"] = shrink_text(slide.get("summary") or "", 88)
        slide["image_caption"] = shrink_text(slide.get("image_caption") or "", 32)
        slide["highlight"] = clamp_text(slide.get("highlight") or "", 78)
        slide["closing"] = shrink_text(slide.get("closing") or "", 84)
        slide["speaker_note"] = clamp_text(slide.get("speaker_note") or "", 220)

        if st == "bullet":
            bullets = [shrink_text(x, 58) for x in (slide.get("bullets") or []) if str(x).strip()]
            if len(bullets) > 4:
                slide_a = dict(slide)
                slide_b = dict(slide)
                slide_a["bullets"] = bullets[:4]
                slide_b["bullets"] = bullets[4:8]
                slide_b["title"] = clamp_text((slide.get("title") or "内容补充") + "（续）", 36)
                new_slides.extend([slide_a, slide_b])
                diagnostics.append(f"拆分页: {slide.get('title') or 'bullet'}")
                continue
            slide["bullets"] = bullets[:4]

        elif st == "two_column":
            left = [shrink_text(x, 48) for x in (slide.get("left_points") or slide.get("bullets")[:2]) if str(x).strip()][:3]
            right = [shrink_text(x, 48) for x in (slide.get("right_points") or slide.get("bullets")[2:]) if str(x).strip()][:3]
            slide["left_points"] = left
            slide["right_points"] = right
            slide["left_title"] = clamp_text(slide.get("left_title") or "左侧要点", 20)
            slide["right_title"] = clamp_text(slide.get("right_title") or "右侧解释", 20)

        elif st == "compare":
            slide["left_points"] = [shrink_text(x, 44) for x in (slide.get("left_points") or []) if str(x).strip()][:4]
            slide["right_points"] = [shrink_text(x, 44) for x in (slide.get("right_points") or []) if str(x).strip()][:4]
            slide["left_title"] = clamp_text(slide.get("left_title") or "方案 A", 18)
            slide["right_title"] = clamp_text(slide.get("right_title") or "方案 B", 18)

        elif st in {"image_left", "image_right"}:
            slide["summary"] = shrink_text(slide.get("summary") or "", 88)
            slide["bullets"] = [shrink_text(x, 52) for x in (slide.get("bullets") or []) if str(x).strip()][:4]
            slide["image_caption"] = shrink_text(slide.get("image_caption") or slide.get("summary") or "页面视觉图", 48)

        elif st == "highlight":
            slide["highlight"] = clamp_text(slide.get("highlight") or slide.get("summary") or slide.get("title") or "核心观点", 68)
            slide["summary"] = shrink_text(slide.get("summary") or "", 88)
            slide["closing"] = shrink_text(slide.get("closing") or slide.get("speaker_note") or "", 78)

        new_slides.append(slide)

    ppt_data["slides"] = new_slides[:10]
    ppt_data["layout_diagnostics"] = diagnostics
    return ppt_data


def is_image_file_usable(image_path: str) -> bool:
    if not image_path or not os.path.exists(image_path):
        return False
    if Image is None:
        return True
    try:
        with Image.open(image_path) as img:
            img.verify()
        with Image.open(image_path) as img:
            w, h = img.size
        if w < 300 or h < 180:
            return False
        ratio = w / h if h else 1.0
        return 0.45 <= ratio <= 2.5
    except Exception:
        return False


def add_card(slide, left, top, width, height, fill_color, line_color=None, line_width=1.0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if line_color is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = line_color
        card.line.width = Pt(line_width)
    return card


def add_footer(slide, theme_tokens, page_text: str):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SAFE_MARGIN_IN), Inches(FOOTER_Y_IN), Inches(SLIDE_W_IN - SAFE_MARGIN_IN*2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = theme_tokens["line"]
    line.line.fill.background()

    page_box, page_tf = safe_add_textbox(slide, 11.45, FOOTER_Y_IN + 0.05, 1.0, 0.24)
    p = page_tf.paragraphs[0]
    p.text = page_text
    style_paragraph(p, size=11, color=theme_tokens["muted"], align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, note_text: str):
    if not note_text:
        return
    try:
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = note_text
    except Exception:
        pass

def render_cover_slide(prs: Presentation, ppt_data: dict, theme_tokens: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0), Inches(0), Inches(0.24), prs.slide_height, theme_tokens["primary"])
    add_full_rect(slide, Inches(9.95), Inches(0), Inches(2.95), Inches(0.95), theme_tokens["accent"])

    title_text, title_size = fit_single_text_to_box(
        ppt_data.get("title", "Selena Presentation"), 8.9, 1.15, 28, 20, max_len=30, bold=True
    )
    _, tf = safe_add_textbox(slide, 1.0, 1.35, 8.9, 1.15, auto_fit=True)
    p = tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])

    subtitle_text, subtitle_size = fit_single_text_to_box(
        ppt_data.get("subtitle", "面向管理层的分析与行动建议"), 8.9, 1.42, 17, 12, max_len=118
    )
    _, sub_tf = safe_add_textbox(slide, 1.02, 2.65, 8.9, 1.42, auto_fit=True)
    sub = sub_tf.paragraphs[0]
    sub.text = subtitle_text
    style_paragraph(sub, size=subtitle_size, color=theme_tokens["muted"])

    _, tag_tf = safe_add_textbox(slide, 1.02, 5.9, 4.0, 0.42, auto_fit=True)
    t = tag_tf.paragraphs[0]
    t.text = f"Theme · {ppt_data.get('theme', 'selena').title()}"
    style_paragraph(t, size=12, bold=True, color=theme_tokens["secondary"])
    return slide
def render_section_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0.85), Inches(1.15), Inches(11.2), Inches(4.0), theme_tokens["accent"])

    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Section", 9.2, 0.88, 28, 20, max_len=34, bold=True)
    _, tf = safe_add_textbox(slide, 1.25, 1.88, 9.2, 0.88, auto_fit=True)
    p = tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])

    summary_text, summary_size = fit_single_text_to_box(
        slide_data.get("summary") or "", 8.35, 2.35, 16, 11, max_len=240
    )
    _, sum_tf = safe_add_textbox(slide, 1.55, 2.85, 8.35, 2.35, auto_fit=True)
    sp = sum_tf.paragraphs[0]
    sp.text = summary_text
    style_paragraph(sp, size=summary_size, color=theme_tokens["muted"], align=PP_ALIGN.CENTER)

    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or slide_data.get("summary") or "")
    return slide
def render_bullet_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0.72), Inches(0.72), Inches(0.14), Inches(0.72), theme_tokens["primary"])

    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Untitled", 10.1, 0.7, 24, 18, max_len=34, bold=True)
    _, title_tf = safe_add_textbox(slide, 1.0, 0.68, 10.1, 0.7, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])

    summary_text, summary_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.0, 0.62, 15, 11, max_len=68)
    _, summary_tf = safe_add_textbox(slide, 1.0, 1.45, 10.0, 0.62, auto_fit=True)
    p2 = summary_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=summary_size, color=theme_tokens["muted"])

    add_card(slide, Inches(0.92), Inches(2.2), Inches(10.9), Inches(3.95), theme_tokens["accent"])
    _, content_tf = safe_add_textbox(slide, 1.25, 2.48, 10.1, 3.35, auto_fit=True)
    bullet_points, bullet_size, _ = fit_bullets_to_box(slide_data.get("bullets"), 10.1, 3.35, preferred_size=17, min_size=11, max_items=4, max_len=42, fallback_text=summary_text)
    populate_bullets(content_tf, bullet_points, bullet_size, theme_tokens["primary"], space_after=5 if bullet_size <= 13 else 7)

    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide
def render_highlight_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(0.72), theme_tokens["accent"])

    small_title, small_size = fit_single_text_to_box(slide_data.get("title") or "Key Point", 10.2, 0.38, 15, 11, max_len=28, bold=True)
    _, small_tf = safe_add_textbox(slide, 0.95, 1.0, 10.2, 0.38, auto_fit=True)
    p0 = small_tf.paragraphs[0]
    p0.text = small_title
    style_paragraph(p0, size=small_size, bold=True, color=theme_tokens["secondary"])

    highlight_text, hl_size = fit_single_text_to_box(slide_data.get("highlight") or slide_data.get("summary") or slide_data.get("title") or "核心观点", 10.4, 1.72, 27, 18, max_len=50, bold=True)
    _, big_tf = safe_add_textbox(slide, 0.95, 1.55, 10.4, 1.72, auto_fit=True)
    p1 = big_tf.paragraphs[0]
    p1.text = highlight_text
    style_paragraph(p1, size=hl_size, bold=True, color=theme_tokens["primary"])

    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.0, 0.88, 16, 11, max_len=70)
    _, summary_tf = safe_add_textbox(slide, 0.95, 3.6, 10.0, 0.88, auto_fit=True)
    p2 = summary_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])

    closing_text, close_size = fit_single_text_to_box(slide_data.get("closing") or slide_data.get("speaker_note") or "", 9.2, 0.64, 13, 10, max_len=62)
    if closing_text:
        _, closing_tf = safe_add_textbox(slide, 0.95, 5.15, 9.2, 0.64, auto_fit=True)
        p3 = closing_tf.paragraphs[0]
        p3.text = closing_text
        style_paragraph(p3, size=close_size, color=theme_tokens["secondary"])

    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide
def render_two_column_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])

    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Two Column", 10.3, 0.7, 23, 18, max_len=34, bold=True)
    _, title_tf = safe_add_textbox(slide, 0.95, 0.68, 10.3, 0.7, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])

    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.1, 0.58, 15, 11, max_len=64)
    _, sum_tf = safe_add_textbox(slide, 0.95, 1.42, 10.1, 0.58, auto_fit=True)
    p2 = sum_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])

    add_card(slide, Inches(0.9), Inches(2.18), Inches(5.1), Inches(3.95), theme_tokens["accent"])
    add_card(slide, Inches(6.25), Inches(2.18), Inches(5.1), Inches(3.95), theme_tokens["accent"])

    left_title, left_title_size = fit_single_text_to_box(slide_data.get("left_title") or "左侧要点", 4.3, 0.34, 16, 12, max_len=18, bold=True)
    left_points, left_size, _ = fit_bullets_to_box(slide_data.get("left_points") or (slide_data.get("bullets") or [])[:2], 4.3, 2.85, preferred_size=15, min_size=11, max_items=3, max_len=34, fallback_text=summary_text)
    _, left_tf = safe_add_textbox(slide, 1.15, 2.45, 4.35, 3.35, auto_fit=True)
    lp = left_tf.paragraphs[0]
    lp.text = left_title
    style_paragraph(lp, size=left_title_size, bold=True, color=theme_tokens["primary"])
    for point in left_points:
        pp = left_tf.add_paragraph()
        pp.text = f"• {point}"
        style_paragraph(pp, size=left_size, color=theme_tokens["primary"])
        pp.space_after = Pt(5)

    right_title, right_title_size = fit_single_text_to_box(slide_data.get("right_title") or "右侧解释", 4.3, 0.34, 16, 12, max_len=18, bold=True)
    right_points, right_size, _ = fit_bullets_to_box(slide_data.get("right_points") or (slide_data.get("bullets") or [])[2:], 4.3, 2.85, preferred_size=15, min_size=11, max_items=3, max_len=34, fallback_text=summary_text)
    _, right_tf = safe_add_textbox(slide, 6.5, 2.45, 4.35, 3.35, auto_fit=True)
    rp = right_tf.paragraphs[0]
    rp.text = right_title
    style_paragraph(rp, size=right_title_size, bold=True, color=theme_tokens["primary"])
    for point in right_points:
        pp = right_tf.add_paragraph()
        pp.text = f"• {point}"
        style_paragraph(pp, size=right_size, color=theme_tokens["primary"])
        pp.space_after = Pt(5)

    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide
def render_image_slide(prs, slide_data, theme_tokens, page_text, image_on_left=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])

    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Visual Slide", 10.6, 0.62, 22, 17, max_len=28, bold=True)
    _, title_tf = safe_add_textbox(slide, 0.82, 0.56, 10.6, 0.62, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])

    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.6, 0.5, 14, 11, max_len=72)
    _, summary_tf = safe_add_textbox(slide, 0.82, 1.22, 10.6, 0.5, auto_fit=True)
    p2 = summary_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])

    img_w = 5.05
    text_w = 4.2
    block_h = 3.75
    img_left = 0.82 if image_on_left else 7.25
    text_left = 6.15 if image_on_left else 0.82
    block_top = 2.0

    add_card(slide, Inches(img_left), Inches(block_top), Inches(img_w), Inches(block_h), theme_tokens["accent"], theme_tokens["line"], 1.0)
    image_path = slide_data.get("image_path") or ""
    if image_path and os.path.exists(image_path):
        add_picture_fit(slide, image_path, Inches(img_left), Inches(block_top), Inches(img_w), Inches(block_h), padding=Inches(0.10), mode="contain")
    else:
        _, placeholder_tf = safe_add_textbox(slide, img_left + 0.3, block_top + 1.45, img_w - 0.6, 0.45, auto_fit=True)
        pp = placeholder_tf.paragraphs[0]
        pp.text = clamp_text(slide_data.get("image_caption") or "视觉图像区域", 16)
        style_paragraph(pp, size=13, bold=True, color=theme_tokens["secondary"], align=PP_ALIGN.CENTER)

    add_card(slide, Inches(text_left), Inches(block_top), Inches(text_w), Inches(block_h), theme_tokens["accent"], theme_tokens["line"], 1.0)
    raw_cap = (slide_data.get("image_caption") or "").strip()
    cap_for_box = raw_cap or (slide_data.get("summary") or "").strip() or "图片与观点保持一致"
    caption_text, cap_size = fit_single_text_to_box(cap_for_box, 3.55, 0.78, 13, 10, max_len=48, bold=True)
    _, text_tf = safe_add_textbox(slide, text_left + 0.22, block_top + 0.2, 3.55, 2.95, auto_fit=True)
    caption = text_tf.paragraphs[0]
    caption.text = caption_text
    style_paragraph(caption, size=cap_size, bold=True, color=theme_tokens["secondary"])
    caption.space_after = Pt(4)

    ref_keys = {k for k in (
        _slide_body_dedupe_key(slide_data.get("summary") or ""),
        _slide_body_dedupe_key(slide_data.get("image_caption") or ""),
        _slide_body_dedupe_key(slide_data.get("title") or ""),
    ) if k}
    _bullets_in = []
    _seen_b = set(ref_keys)
    for b in slide_data.get("bullets") or []:
        bk = _slide_body_dedupe_key(str(b))
        if bk and bk not in _seen_b:
            _seen_b.add(bk)
            _bullets_in.append(b)
    image_points, image_bullet_size, _ = fit_bullets_to_box(
        _bullets_in, 3.55, 2.15, preferred_size=12, min_size=10, max_items=4, max_len=42, fallback_text=summary_text
    )
    for point in image_points:
        pt = (point or "").strip()
        if not pt:
            continue
        pp = text_tf.add_paragraph()
        pp.text = f"• {pt}"
        style_paragraph(pp, size=image_bullet_size, color=theme_tokens["primary"])
        pp.space_after = Pt(3)

    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide
def render_compare_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])

    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Compare", 10.4, 0.7, 23, 18, max_len=34, bold=True)
    _, title_tf = safe_add_textbox(slide, 0.95, 0.68, 10.4, 0.7, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])

    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.2, 0.56, 15, 11, max_len=60)
    _, sub_tf = safe_add_textbox(slide, 0.95, 1.42, 10.2, 0.56, auto_fit=True)
    p2 = sub_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])

    add_full_rect(slide, Inches(0.95), Inches(2.08), Inches(4.95), Inches(0.46), theme_tokens["primary"])
    add_full_rect(slide, Inches(6.28), Inches(2.08), Inches(4.95), Inches(0.46), theme_tokens["secondary"])

    _, l_head_tf = safe_add_textbox(slide, 1.12, 2.15, 4.55, 0.26, auto_fit=True)
    lp = l_head_tf.paragraphs[0]
    lp.text = clamp_text(slide_data.get("left_title") or "方案 A", 14)
    style_paragraph(lp, size=13, bold=True, color=theme_tokens["light"])

    _, r_head_tf = safe_add_textbox(slide, 6.45, 2.15, 4.55, 0.26, auto_fit=True)
    rp = r_head_tf.paragraphs[0]
    rp.text = clamp_text(slide_data.get("right_title") or "方案 B", 14)
    style_paragraph(rp, size=13, bold=True, color=theme_tokens["light"])

    add_card(slide, Inches(0.95), Inches(2.65), Inches(4.95), Inches(3.55), theme_tokens["accent"])
    add_card(slide, Inches(6.28), Inches(2.65), Inches(4.95), Inches(3.55), theme_tokens["accent"])

    left_points, left_size, _ = fit_bullets_to_box(slide_data.get("left_points"), 4.2, 2.9, preferred_size=14, min_size=10, max_items=4, max_len=30, fallback_text=summary_text)
    _, left_tf = safe_add_textbox(slide, 1.22, 2.95, 4.18, 2.95, auto_fit=True)
    populate_bullets(left_tf, left_points, left_size, theme_tokens["primary"], space_after=4)

    right_points, right_size, _ = fit_bullets_to_box(slide_data.get("right_points"), 4.2, 2.9, preferred_size=14, min_size=10, max_items=4, max_len=30, fallback_text=summary_text)
    _, right_tf = safe_add_textbox(slide, 6.55, 2.95, 4.18, 2.95, auto_fit=True)
    populate_bullets(right_tf, right_points, right_size, theme_tokens["primary"], space_after=4)

    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide
def build_pptx_file(ppt_data: dict) -> BytesIO:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed.")
    budget_fn = modular_enforce_slide_content_budget or enforce_slide_content_budget
    ppt_data = budget_fn(ppt_data)
    try:
        from ppt.layout_engine import finalize_slide_layouts
        finalize_slide_layouts(ppt_data.get("slides", []))
    except Exception:
        pass
    ppt_data = ensure_slide_images(ppt_data, topic=ppt_data.get("title") or "Presentation", max_images=2)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    theme_tokens = get_theme_tokens(ppt_data.get("theme", "selena"))
    render_cover_slide(prs, ppt_data, theme_tokens)

    for idx, slide_data in enumerate(ppt_data.get("slides", []), start=2):
        slide_type = slide_data.get("type") or "bullet"
        page_text = f"{idx - 1}/{len(ppt_data.get('slides', []))}"
        if slide_type == "section":
            render_section_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "highlight":
            render_highlight_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "timeline":
            try:
                from ppt.renderers import render_timeline_slide as _render_timeline_slide
                _render_timeline_slide(prs, slide_data, theme_tokens, page_text)
            except Exception:
                fb = dict(slide_data)
                fb["type"] = "bullet"
                fb["bullets"] = fb.get("timeline_steps") or fb.get("bullets") or []
                render_bullet_slide(prs, fb, theme_tokens, page_text)
        elif slide_type == "two_column":
            render_two_column_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "compare":
            render_compare_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "image_left":
            render_image_slide(prs, slide_data, theme_tokens, page_text, image_on_left=True)
        elif slide_type == "image_right":
            render_image_slide(prs, slide_data, theme_tokens, page_text, image_on_left=False)
        else:
            render_bullet_slide(prs, slide_data, theme_tokens, page_text)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def summarize_ppt_for_chat(ppt_data: dict) -> str:
    slides = ppt_data.get("slides") or []
    overview = "\n".join([f"- {s.get('title', '未命名页面')}（{s.get('type', 'bullet')}）" for s in slides[:6]])
    image_count = int(ppt_data.get("generated_image_count") or 0)
    image_line = f"**配图**：已自动生成并插入 {image_count} 张页面图片\n" if image_count else "**配图**：这次未插入图片，已按文字版式生成\n"
    return (
        f"我已经按 **{ppt_data.get('theme', 'selena').title()}** 风格生成了一份可下载的 PPT。\n\n"
        f"**标题**：{ppt_data.get('title', '未命名演示')}\n"
        f"**副标题**：{ppt_data.get('subtitle', '')}\n"
        f"**页数**：{len(slides) + 1} 页（含封面）\n"
        f"{image_line}\n"
        f"**目录预览**：\n{overview}\n\n"
        f"这版会根据主题自动挑选 1 到 2 页做视觉化表达，并把生成的图片直接嵌入到 PPT 版式里。"
    )

def save_generated_ppt(ppt_bytes: BytesIO, filename: str, conversation_id: str = "default", ppt_title: str = "", slide_count: int = 0, image_count: int = 0) -> str:
    file_id = "ppt_" + uuid.uuid4().hex
    temp_dir = tempfile.gettempdir()
    path = os.path.join(temp_dir, f"{file_id}.pptx")
    with open(path, "wb") as f:
        f.write(ppt_bytes.getvalue())
    generated_ppts[file_id] = {
        "path": path,
        "filename": filename,
        "ppt_title": ppt_title or filename,
        "slide_count": int(slide_count or 0),
        "image_count": int(image_count or 0),
        "conversation_id": conversation_id,
        "created_at": datetime.utcnow().isoformat() + "Z",
    }
    _persist_generated_ppts()
    _attach_ppt_to_conversation(conversation_id, file_id)
    return file_id



# =========================
# PPT service facade
# This layer is the bridge between app.py and the future ppt/ package.
# Routes should call these functions instead of directly touching low-level
# PPT functions, so later we can move code out file by file without breaking
# the API layer.
# =========================
def ppt_generate_structure_service(topic: str, upload_id: str = "") -> dict:
    generator_fn = modular_generate_ppt_structure or generate_ppt_structure
    data = generator_fn(topic, upload_id)
    if not isinstance(data, dict):
        raise ValueError("PPT generator returned invalid data.")
    return data


def ppt_normalize_service(ppt_data: dict) -> dict:
    normalizer_fn = modular_normalize_ppt_json or normalize_ppt_json
    budget_fn = modular_enforce_slide_content_budget or enforce_slide_content_budget
    ppt_data = normalizer_fn(ppt_data)
    ppt_data = budget_fn(ppt_data)
    return ppt_data


def ppt_build_file_service(ppt_data: dict) -> BytesIO:
    builder_fn = modular_build_pptx_file or build_pptx_file
    return builder_fn(ppt_data)


def ppt_summarize_service(ppt_data: dict) -> str:
    summary_fn = modular_summarize_ppt_for_chat or summarize_ppt_for_chat
    return summary_fn(ppt_data)


def create_ppt_package(topic: str, upload_id: str = "") -> tuple[dict, BytesIO]:
    ppt_data = ppt_generate_structure_service(topic, upload_id)
    ppt_data = ppt_normalize_service(ppt_data)
    ppt_file = ppt_build_file_service(ppt_data)
    return ppt_data, ppt_file


# =========================
# Routes
# =========================
@app.route("/")
def home():
    return render_template(
        "index.html",
        profile=ASSISTANT_PROFILE,
        modes=ASSISTANT_PROFILE.get("modes", {}),
        image_defaults=ASSISTANT_PROFILE.get("image_generation", {})
    )


@app.route("/upload-file", methods=["POST"])
def upload_file():
    conversation_id = (request.form.get("conversation_id") or "default").strip()
    file = request.files.get("file")
    if not file or not file.filename:
        return jsonify({"error": "No file uploaded."}), 400

    filename = file.filename.strip()
    extension = os.path.splitext(filename)[1].lower()

    if extension not in ALLOWED_EXTENSIONS:
        return jsonify({"error": "Unsupported file type.", "allowed_extensions": sorted(ALLOWED_EXTENSIONS)}), 400

    try:
        raw_bytes = file.read()
        text = read_uploaded_file(filename, raw_bytes)
    except Exception as e:
        print("File read error:", e)
        return jsonify({"error": f"Failed to read file content: {str(e)}"}), 400

    truncated = len(text) > MAX_FILE_CHARS
    content = text[:MAX_FILE_CHARS]
    upload_id = "upload_" + uuid.uuid4().hex

    uploaded_files[upload_id] = {"filename": filename, "extension": extension, "content": content, "created_at": datetime.utcnow().isoformat() + "Z", "conversation_id": conversation_id}
    _persist_uploaded_files()
    _attach_upload_to_conversation(conversation_id, upload_id)
    return jsonify({
        "upload_id": upload_id,
        "filename": filename,
        "size_chars": len(content),
        "truncated": truncated,
        "suggested_mode": infer_mode_by_filename(filename)
    })


@app.route("/download-ppt/<file_id>")
def download_ppt(file_id):
    info = generated_ppts.get(file_id)
    if not info or not os.path.exists(info["path"]):
        return jsonify({"error": "PPT file not found or expired."}), 404
    return send_file(
        info["path"],
        as_attachment=True,
        download_name=info["filename"],
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@app.route("/list-ppts", methods=["GET"])
def list_ppts():
    conversation_id = (request.args.get("conversation_id") or "").strip()
    items = []
    for file_id, info in generated_ppts.items():
        if conversation_id and info.get("conversation_id") != conversation_id:
            continue
        items.append({
            "file_id": file_id,
            "filename": info.get("filename"),
            "ppt_title": info.get("ppt_title") or info.get("filename"),
            "slide_count": int(info.get("slide_count") or 0),
            "image_count": int(info.get("image_count") or 0),
            "created_at": info.get("created_at"),
            "download_url": f"/download-ppt/{file_id}",
        })
    items.sort(key=lambda x: x.get("created_at") or "", reverse=True)
    return jsonify({"items": items})


@app.route("/delete-ppt/<file_id>", methods=["DELETE"])
def delete_ppt(file_id):
    info = generated_ppts.get(file_id)
    if not info:
        return jsonify({"error": "PPT file not found."}), 404

    conversation_id = info.get("conversation_id") or ""
    if conversation_id and conversation_id in conversation_resources:
        try:
            conversation_resources[conversation_id]["ppts"] = [
                fid for fid in conversation_resources[conversation_id].get("ppts", []) if fid != file_id
            ]
        except Exception:
            pass

    ok = delete_generated_ppt_file(file_id)
    if not ok:
        return jsonify({"error": "PPT file could not be deleted."}), 500
    return jsonify({"status": "deleted", "file_id": file_id})


@app.route("/chat-ppt", methods=["POST"])
def chat_ppt():
    data = request.get_json(silent=True) or {}
    user_message = (data.get("message") or "").strip()
    conversation_id = (data.get("conversation_id") or "default").strip()
    upload_id = (data.get("upload_id") or "").strip()

    if not user_message and not upload_id:
        return jsonify({"error": "Message or uploaded file is required."}), 400

    uploaded = uploaded_files.get(upload_id) if upload_id else None
    topic = user_message or (uploaded["filename"] if uploaded else "Selena Presentation")

    try:
        ppt_data, ppt_bytes = create_ppt_package(topic, upload_id)
        filename = safe_filename(ppt_data.get("title") or topic or "Selena_Presentation") + ".pptx"
        file_id = save_generated_ppt(
            ppt_bytes,
            filename,
            conversation_id=conversation_id,
            ppt_title=ppt_data.get("title", filename),
            slide_count=len(ppt_data.get("slides", [])) + 1,
            image_count=int(ppt_data.get("generated_image_count") or 0),
        )
        reply_text = ppt_summarize_service(ppt_data)
        download_url = f"/download-ppt/{file_id}"

        history = get_history(conversation_id)
        visible_user = build_file_augmented_prompt(uploaded, user_message, "ppt") if uploaded else user_message
        history.append({"role": "user", "content": visible_user})
        history.append({
            "role": "assistant",
            "content": reply_text,
            "type": "ppt",
            "download_url": download_url,
            "filename": filename,
            "ppt_title": ppt_data.get("title", filename),
            "slide_count": len(ppt_data.get("slides", [])) + 1,
            "image_count": int(ppt_data.get("generated_image_count") or 0),
            "file_id": file_id,
            "created_at": generated_ppts.get(file_id, {}).get("created_at"),
        })

        return jsonify({
            "text": reply_text,
            "download_url": download_url,
            "filename": filename,
            "ppt_title": ppt_data.get("title", filename),
            "slide_count": len(ppt_data.get("slides", [])) + 1,
            "image_count": int(ppt_data.get("generated_image_count") or 0),
            "file_id": file_id,
            "created_at": generated_ppts.get(file_id, {}).get("created_at"),
        })
    except Exception as e:
        print("PPT chat generation error:", e)
        return jsonify({"error": f"PPT generation failed: {str(e)}"}), 500


@app.route("/chat-stream", methods=["POST"])
def chat_stream():
    data = request.get_json(silent=True) or {}
    user_message = (data.get("message") or "").strip()
    conversation_id = (data.get("conversation_id") or "default").strip()
    mode = (data.get("mode") or "chat").strip()
    upload_id = (data.get("upload_id") or "").strip()

    uploaded = uploaded_files.get(upload_id) if upload_id else None
    if not user_message and not uploaded:
        return Response(ASSISTANT_PROFILE["empty_message"], mimetype="text/plain; charset=utf-8", status=400)

    model_input = build_file_augmented_prompt(uploaded, user_message, mode) if uploaded else user_message
    chat_history = get_history(conversation_id)
    chat_history.append({"role": "user", "content": model_input})
    system_prompt = build_system_prompt(mode)

    @stream_with_context
    def generate():
        full_reply = ""
        try:
            stream = client.chat.completions.create(
                model="openai/gpt-4o-mini",
                messages=[{"role": "system", "content": system_prompt}] + chat_history,
                stream=True
            )
            for chunk in stream:
                delta = ""
                try:
                    delta = chunk.choices[0].delta.content or ""
                except Exception:
                    delta = ""
                if delta:
                    full_reply += delta
                    yield delta
            if full_reply.strip():
                chat_history.append({"role": "assistant", "content": full_reply})
        except GeneratorExit:
            if full_reply.strip():
                chat_history.append({"role": "assistant", "content": full_reply})
            raise
        except Exception as e:
            print("Streaming chat error:", e)
            if chat_history and chat_history[-1]["role"] == "user":
                chat_history.pop()
            yield ASSISTANT_PROFILE["server_error"]

    return Response(generate(), mimetype="text/plain; charset=utf-8")


@app.route("/generate-image", methods=["POST"])
def generate_image():
    data = request.get_json(silent=True) or {}
    prompt = (data.get("prompt") or "").strip()
    conversation_id = (data.get("conversation_id") or "default").strip()
    aspect_ratio = (data.get("aspect_ratio") or ASSISTANT_PROFILE["image_generation"]["aspect_ratio"]).strip()
    image_size = (data.get("image_size") or ASSISTANT_PROFILE["image_generation"]["image_size"]).strip()

    if not prompt:
        return jsonify({"error": "Prompt is required."}), 400

    image_config = ASSISTANT_PROFILE.get("image_generation", {})
    model = image_config.get("model", "google/gemini-2.5-flash-image")

    try:
        response = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={
                "model": model,
                "messages": [{"role": "user", "content": prompt}],
                "modalities": ["image", "text"],
                "stream": False,
                "image_config": {"aspect_ratio": aspect_ratio, "image_size": image_size}
            },
            timeout=180
        )
        response.raise_for_status()
        result = response.json()

        choices = result.get("choices") or []
        message = ((choices[0] if choices else {}) or {}).get("message") or {}
        text = message.get("content") or "我已经为你准备好了图像。"
        image_urls = extract_image_urls(result)
        if not image_urls:
            print("Image response with no parsed images:", result)
            return jsonify({"error": "The image model returned no images.", "details": "Try switching to another supported image model."}), 502

        chat_history = get_history(conversation_id)
        chat_history.append({"role": "user", "content": prompt})
        chat_history.append({"role": "assistant", "content": text, "type": "image", "images": image_urls})
        return jsonify({"images": image_urls, "text": text})
    except requests.HTTPError as e:
        error_text = e.response.text if e.response is not None else str(e)
        print("Image generation HTTP error:", error_text)
        return jsonify({"error": "Image generation request failed.", "details": error_text[:500]}), 500
    except Exception as e:
        print("Image generation error:", e)
        return jsonify({"error": "Image generation failed."}), 500


@app.route("/title", methods=["POST"])
def generate_title():
    data = request.get_json(silent=True) or {}
    user_message = (data.get("message") or "").strip()
    if not user_message:
        return jsonify({"title": "新聊天"})
    return jsonify({"title": generate_title_from_message(user_message)})


@app.route("/clear", methods=["POST"])
def clear_chat():
    data = request.get_json(silent=True) or {}
    conversation_id = (data.get("conversation_id") or "default").strip()
    clear_files = bool(data.get("clear_files", True))
    conversation_histories[conversation_id] = []
    if clear_files:
        clear_conversation_resources(conversation_id)
    return jsonify({"status": "cleared", "conversation_id": conversation_id, "files_cleared": clear_files})


if __name__ == "__main__":
    app.run(debug=False, threaded=True)
