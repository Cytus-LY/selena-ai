import os
import re
import math
import tempfile
from io import BytesIO

import requests

try:
    from PIL import Image
except Exception:
    Image = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    RGBColor = None
    PP_ALIGN = None
    MSO_AUTO_SIZE = None
    MSO_VERTICAL_ANCHOR = None
    MSO_SHAPE = None

try:
    from .normalizer import enforce_slide_content_budget as normalize_ppt_data
except Exception:
    from normalizer import enforce_slide_content_budget as normalize_ppt_data

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
FOOTER_Y_IN = 6.82


def _walk_json(node):
    if isinstance(node, dict):
        yield node
        for value in node.values():
            yield from _walk_json(value)
    elif isinstance(node, list):
        for item in node:
            yield from _walk_json(item)
    else:
        yield node


def extract_image_urls(result_json: dict) -> list[str]:
    image_urls = []
    for node in _walk_json(result_json):
        if not isinstance(node, dict):
            continue
        direct = node.get("url")
        if isinstance(direct, str) and direct.startswith(("http://", "https://", "data:image/")):
            image_urls.append(direct)
        for key in ("image_url", "imageUrl"):
            nested = node.get(key)
            if isinstance(nested, dict):
                url = nested.get("url")
                if isinstance(url, str) and url.startswith(("http://", "https://", "data:image/")):
                    image_urls.append(url)
    seen = set()
    unique_urls = []
    for url in image_urls:
        if url not in seen:
            seen.add(url)
            unique_urls.append(url)
    return unique_urls


def extract_image_b64_strings(result_json: dict) -> list[str]:
    candidates = []
    b64_keys = {"b64_json", "b64", "base64", "image_base64", "imageBase64", "png", "jpeg", "jpg", "webp"}
    for node in _walk_json(result_json):
        if not isinstance(node, dict):
            continue
        for key, value in node.items():
            if key in b64_keys and isinstance(value, str) and len(value) > 200:
                candidates.append(value.strip())
    deduped = []
    seen = set()
    for item in candidates:
        if item not in seen:
            seen.add(item)
            deduped.append(item)
    return deduped


def save_base64_image(b64_data: str, suffix: str = ".png") -> str:
    import base64
    raw = b64_data.strip()
    if raw.startswith("data:image/"):
        header, raw = raw.split(",", 1)
        if "jpeg" in header or "jpg" in header:
            suffix = ".jpg"
        elif "webp" in header:
            suffix = ".webp"
    fd, path = tempfile.mkstemp(prefix="selena_b64_", suffix=suffix)
    os.close(fd)
    with open(path, "wb") as f:
        f.write(base64.b64decode(raw))
    return path


def download_binary_file(url: str, suffix: str = ".png") -> str:
    resp = requests.get(url, timeout=180, stream=True)
    resp.raise_for_status()
    content_type = (resp.headers.get("Content-Type") or "").lower()
    if not content_type.startswith("image/"):
        raise ValueError(f"URL did not return an image. Content-Type={content_type}")
    fd, path = tempfile.mkstemp(prefix="selena_img_", suffix=suffix)
    os.close(fd)
    with open(path, "wb") as f:
        for chunk in resp.iter_content(8192):
            if chunk:
                f.write(chunk)
    return path


def download_or_decode_image(source: str, suffix: str = ".png") -> str:
    source = (source or "").strip()
    if source.startswith("data:image/"):
        return save_base64_image(source, suffix=suffix)
    ext = os.path.splitext(source.split("?", 1)[0])[1].lower()
    if ext in {".png", ".jpg", ".jpeg", ".webp"}:
        suffix = ext
    return download_binary_file(source, suffix=suffix)


def choose_ppt_image_model() -> str:
    return (
        os.getenv("PPT_IMAGE_MODEL")
        or os.getenv("OPENROUTER_IMAGE_MODEL")
        or "google/gemini-2.5-flash-image"
    ).strip()


def llm_generate_image_prompt(topic: str, slide_data: dict, theme: str) -> str:
    title = (slide_data.get("title") or "").strip()
    summary = (slide_data.get("summary") or "").strip()
    bullets = [str(x).strip() for x in (slide_data.get("bullets") or []) if str(x).strip()]
    caption = (slide_data.get("image_caption") or "").strip()
    parts = [
        "Clean presentation-ready illustration",
        f"theme: {theme}",
        f"topic: {topic}" if topic else "",
        f"slide title: {title}" if title else "",
        f"focus: {summary}" if summary else "",
        f"details: {'; '.join(bullets[:3])}" if bullets else "",
        f"caption hint: {caption}" if caption else "",
        "modern, uncluttered, no text, no watermark, professional lighting",
    ]
    return ", ".join([p for p in parts if p])


def _norm_text(text: str) -> str:
    return " ".join((text or "").replace("•", "").split()).strip().lower()


def _dedupe_preserve_order(items):
    seen = set()
    out = []
    for item in items or []:
        s = str(item or "").strip()
        key = _norm_text(s)
        if not s or key in seen:
            continue
        seen.add(key)
        out.append(s)
    return out


def choose_visual_candidate_score(slide: dict) -> int:
    slide_type = slide.get("type") or "bullet"
    if slide_type not in {"bullet", "highlight", "two_column", "compare"}:
        return -999
    score = 0
    summary = slide.get("summary") or ""
    bullets = slide.get("bullets") or []
    total_points = len(bullets) + len(slide.get("left_points") or []) + len(slide.get("right_points") or [])
    score += 12 if slide_type == "highlight" else 0
    score += 6 if 18 <= len(summary) <= 58 else 0
    score += 5 if total_points <= 3 else 0
    score -= 5 if total_points >= 6 else 0
    score -= 8 if slide_type == "compare" else 0
    score += 4 if slide.get("highlight") else 0
    score -= 9 if slide.get("_continued") else 0
    return score


def ensure_visual_slides(ppt_data: dict, topic: str, min_images: int = 1, max_images: int = 2) -> dict:
    slides = list(ppt_data.get("slides") or [])
    existing = [i for i, s in enumerate(slides) if (s.get("type") or "") in {"image_left", "image_right"}]
    if len(existing) >= min_images:
        return ppt_data
    candidates = []
    for idx, slide in enumerate(slides):
        score = choose_visual_candidate_score(slide)
        if score > -500:
            candidates.append((score, idx))
    candidates.sort(reverse=True)
    target_count = min(max(min_images, 1), max_images)
    chosen = set(existing)
    for _, idx in candidates:
        if len(chosen) >= target_count:
            break
        chosen.add(idx)
    for n, idx in enumerate(sorted(chosen)):
        slide = slides[idx]
        if (slide.get("type") or "") in {"image_left", "image_right"}:
            continue
        slide["type"] = "image_left" if n % 2 == 0 else "image_right"
        slide.setdefault("image_caption", clamp_text(slide.get("summary") or slide.get("title") or topic, 30))
    ppt_data["slides"] = slides
    return ppt_data


def generate_presentation_image_assets(prompt: str, aspect_ratio: str = "16:9", image_size: str = "1K") -> dict:
    api_key = os.getenv("OPENROUTER_API_KEY")
    model = choose_ppt_image_model()
    print("==== IMAGE DEBUG START ====")
    print("prompt exists:", bool(prompt))
    print("prompt preview:", repr((prompt or "")[:120]))
    print("OPENROUTER_API_KEY exists:", bool(api_key))
    print("OPENROUTER_API_KEY prefix:", (api_key[:12] + "...") if api_key else None)
    print("PPT_IMAGE_MODEL:", repr(model))
    print("==== IMAGE DEBUG END ====")
    if not prompt or not api_key or not model:
        print("Image generation skipped: missing prompt, OPENROUTER_API_KEY, or PPT_IMAGE_MODEL")
        return {"urls": [], "b64": [], "raw": {}}
    payload = {
        "model": model,
        "messages": [{
            "role": "user",
            "content": (
                "Generate one clean presentation-style illustration. "
                "No text, no watermark, no logo, no UI, no collage. "
                f"Aspect ratio preference: {aspect_ratio}. "
                f"Prompt: {prompt}"
            ),
        }],
    }
    try:
        response = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json=payload,
            timeout=240,
        )
        if response.status_code != 200:
            print("OpenRouter image request failed:")
            print("Status:", response.status_code)
            try:
                print("Body:", response.text[:1200])
            except Exception:
                pass
            response.raise_for_status()
        result = response.json()
        urls = extract_image_urls(result)
        b64 = extract_image_b64_strings(result)
        print(f"Image API returned urls: {len(urls)} b64: {len(b64)}")
        return {"urls": urls, "b64": b64, "raw": result}
    except Exception as e:
        print("OpenRouter image request exception:", e)
        return {"urls": [], "b64": [], "raw": {}}


def is_image_file_usable(image_path: str) -> bool:
    if not image_path or not os.path.exists(image_path) or Image is None:
        return bool(image_path and os.path.exists(image_path))
    try:
        with Image.open(image_path) as img:
            w, h = img.size
        if w < 320 or h < 240:
            return False
        ratio = max(w / h, h / w)
        if ratio > 3.2:
            return False
        return True
    except Exception:
        return False


def try_generate_slide_image_file(prompt: str) -> str:
    attempts = [("16:9", "1K"), ("4:3", "1K"), ("1:1", "1K")]
    for aspect_ratio, image_size in attempts:
        print(f"Trying image generation: aspect_ratio={aspect_ratio}, image_size={image_size}")
        try:
            assets = generate_presentation_image_assets(prompt=prompt, aspect_ratio=aspect_ratio, image_size=image_size)
            path = ""
            if assets["b64"]:
                print("Using base64 image payload")
                path = save_base64_image(assets["b64"][0], ".png")
            elif assets["urls"]:
                print("Using structured image url source:", str(assets["urls"][0])[:180])
                path = download_or_decode_image(assets["urls"][0], ".png")
            else:
                print("Image API returned no usable image data.")
            if path and is_image_file_usable(path):
                print("Image generation success:", path)
                return path
            elif path:
                print("Generated image file was unusable:", path)
        except Exception as e:
            print(f"Image generation attempt failed ({aspect_ratio}, {image_size}):", e)
    print("All image generation attempts failed.")
    return ""


def ensure_slide_images(ppt_data: dict, topic: str = "", max_images: int = 2) -> dict:
    ppt_data = ensure_visual_slides(ppt_data, topic=topic or ppt_data.get("title", "Presentation"), min_images=1, max_images=max_images)
    generated = 0
    for slide in ppt_data.get("slides", []):
        if generated >= max_images:
            break
        slide_type = slide.get("type") or "bullet"
        prompt = (slide.get("image_prompt") or "").strip()
        if slide_type not in {"image_left", "image_right"}:
            continue
        if not prompt:
            prompt = llm_generate_image_prompt(topic or ppt_data.get("title", "Presentation"), slide, ppt_data.get("theme", "selena"))
            slide["image_prompt"] = prompt
        try:
            image_path = try_generate_slide_image_file(prompt)
            slide["image_path"] = image_path or ""
            if image_path:
                generated += 1
        except Exception as e:
            print("Slide image generation error:", e)
            slide["image_path"] = ""
    ppt_data["generated_image_count"] = generated
    return ppt_data


def add_picture_fit(slide, image_path, left, top, width, height, padding=Inches(0.08), mode="contain"):
    if not image_path or not os.path.exists(image_path):
        return None
    inner_left = max(left + padding, Inches(0.05))
    inner_top = max(top + padding, Inches(0.05))
    max_right = Inches(SLIDE_W_IN - 0.08)
    max_bottom = Inches(SLIDE_H_IN - 0.08)
    inner_width = max(min(width - padding * 2, max_right - inner_left), Inches(0.2))
    inner_height = max(min(height - padding * 2, max_bottom - inner_top), Inches(0.2))
    if Image is None:
        return slide.shapes.add_picture(image_path, inner_left, inner_top, width=inner_width, height=inner_height)
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    if not img_w or not img_h:
        return None
    box_ratio = float(inner_width) / float(inner_height) if float(inner_height) else 1.0
    img_ratio = img_w / img_h
    mode = (mode or "contain").lower()
    if mode == "cover":
        if img_ratio > box_ratio:
            pic_height = inner_height
            pic_width = min(inner_height * img_ratio, Inches(SLIDE_W_IN - 0.16))
            pic_left = max(Inches(0.05), inner_left - (pic_width - inner_width) / 2)
            pic_top = inner_top
        else:
            pic_width = inner_width
            pic_height = min(inner_width / img_ratio, Inches(SLIDE_H_IN - 0.16))
            pic_left = inner_left
            pic_top = max(Inches(0.05), inner_top - (pic_height - inner_height) / 2)
    else:
        if img_ratio > box_ratio:
            pic_width = inner_width
            pic_height = inner_width / img_ratio if img_ratio else inner_height
            pic_left = inner_left
            pic_top = inner_top + (inner_height - pic_height) / 2
        else:
            pic_height = inner_height
            pic_width = inner_height * img_ratio
            pic_left = inner_left + (inner_width - pic_width) / 2
            pic_top = inner_top
    pic_width = min(pic_width, max_right - pic_left)
    pic_height = min(pic_height, max_bottom - pic_top)
    return slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width, height=pic_height)


def get_theme_tokens(theme: str) -> dict:
    theme = (theme or "selena").lower()
    if theme == "business":
        return {
            "primary": RGBColor(15, 23, 42),
            "secondary": RGBColor(37, 99, 235),
            "accent": RGBColor(219, 234, 254),
            "muted": RGBColor(71, 85, 105),
            "light": RGBColor(248, 250, 252),
            "line": RGBColor(191, 219, 254),
        }
    if theme == "academic":
        return {
            "primary": RGBColor(30, 64, 175),
            "secondary": RGBColor(14, 116, 144),
            "accent": RGBColor(224, 242, 254),
            "muted": RGBColor(71, 85, 105),
            "light": RGBColor(255, 255, 255),
            "line": RGBColor(186, 230, 253),
        }
    return {
        "primary": RGBColor(79, 70, 229),
        "secondary": RGBColor(139, 92, 246),
        "accent": RGBColor(237, 233, 254),
        "muted": RGBColor(71, 85, 105),
        "light": RGBColor(255, 255, 255),
        "line": RGBColor(196, 181, 253),
    }


def add_full_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", auto_fit=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_fit else MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    if text:
        tf.text = text
    return box, tf


def style_paragraph(paragraph, size=20, bold=False, color=None, align=None):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def clamp_text(text: str, max_len: int = 120) -> str:
    text = (text or "").strip()
    if len(text) <= max_len:
        return text
    return text[: max_len - 1].rstrip() + "…"


def adaptive_font_size(text: str, default: int, minimum: int, breakpoints=None) -> int:
    text = (text or "").strip()
    breakpoints = breakpoints or []
    size = default
    for char_count, candidate in sorted(breakpoints):
        if len(text) > char_count:
            size = min(size, candidate)
    return max(minimum, size)


def detect_text_language(text: str) -> str:
    text = text or ""
    zh_count = sum(1 for ch in text if "一" <= ch <= "鿿")
    return "zh" if zh_count >= max(2, len(text) * 0.25) else "en"


def estimate_lines(text: str, box_width_in: float, font_size_pt: int, prefix_chars: float = 0.0) -> int:
    text = (text or "").strip()
    if not text:
        return 1
    lang = detect_text_language(text)
    avg_char_factor = 0.95 if lang == "zh" else 0.56
    box_width_pt = max(24.0, box_width_in * 72.0)
    chars_per_line = max(1.0, box_width_pt / (font_size_pt * avg_char_factor))
    effective_len = len(text) + prefix_chars
    return max(1, int(math.ceil(effective_len / chars_per_line)))


def estimate_text_block_height(texts, box_width_in: float, font_size_pt: int, line_spacing=1.22, paragraph_gap_pt=6, prefix_chars: float = 0.0) -> float:
    total_pt = 0.0
    items = list(texts or [])
    if not items:
        return font_size_pt * line_spacing / 72.0
    for idx, text in enumerate(items):
        lines = estimate_lines(text, box_width_in, font_size_pt, prefix_chars=prefix_chars)
        total_pt += lines * font_size_pt * line_spacing
        if idx < len(items) - 1:
            total_pt += paragraph_gap_pt
    return total_pt / 72.0


def shrink_text(text: str, max_len: int) -> str:
    text = re.sub(r"\s+", " ", (text or "").strip())
    text = re.sub(r"^(首先|其次|另外|此外|总的来说|总体来看|需要注意的是)[，,:： ]*", "", text)
    text = re.sub(r"(可以看到|我们可以看到|这说明|这意味着)[，,:： ]*", "", text)
    return clamp_text(text, max_len)


def fit_bullets_to_box(points, box_width_in, box_height_in, preferred_size=18, min_size=14, max_items=4, max_len=54, fallback_text=""):
    cleaned = []
    for item in points or []:
        txt = str(item).strip()
        if txt:
            cleaned.append(txt)
    if not cleaned and fallback_text:
        cleaned = [fallback_text]
    cleaned = [shrink_text(x, max_len) for x in cleaned[:max_items]]
    if not cleaned:
        return [], preferred_size, False
    for size in range(preferred_size, min_size - 1, -1):
        needed_h = estimate_text_block_height(cleaned, box_width_in, size, line_spacing=1.18, paragraph_gap_pt=7, prefix_chars=2.2)
        if needed_h <= box_height_in:
            return cleaned, size, False
    trimmed = cleaned[:]
    reduced = False
    while len(trimmed) > 2:
        trimmed = trimmed[:-1]
        reduced = True
        for size in range(preferred_size, min_size - 1, -1):
            needed_h = estimate_text_block_height(trimmed, box_width_in, size, line_spacing=1.18, paragraph_gap_pt=7, prefix_chars=2.2)
            if needed_h <= box_height_in:
                return trimmed, size, reduced
    compact = [shrink_text(x, max(16, max_len - 10)) for x in trimmed]
    for size in range(preferred_size, min_size - 1, -1):
        needed_h = estimate_text_block_height(compact, box_width_in, size, line_spacing=1.16, paragraph_gap_pt=5, prefix_chars=2.2)
        if needed_h <= box_height_in:
            return compact, size, True
    return compact[:2], min_size, True


def fit_single_text_to_box(text: str, box_width_in: float, box_height_in: float, preferred_size: int, min_size: int, max_len: int = None, bold: bool = False):
    txt = (text or "").strip()
    if max_len:
        txt = shrink_text(txt, max_len)
    if not txt:
        return "", preferred_size
    for size in range(preferred_size, min_size - 1, -1):
        needed_h = estimate_text_block_height([txt], box_width_in, size, line_spacing=1.15 if bold else 1.18, paragraph_gap_pt=0)
        if needed_h <= box_height_in:
            return txt, size
    if max_len:
        for cap in range(max(min(max_len, len(txt)), min_size * 2), max(10, min(max_len, len(txt)) // 2), -4):
            shorter = shrink_text(txt, cap)
            needed_h = estimate_text_block_height([shorter], box_width_in, min_size, line_spacing=1.15 if bold else 1.18, paragraph_gap_pt=0)
            if needed_h <= box_height_in:
                return shorter, min_size
    return shrink_text(txt, max_len or max(18, int(box_width_in * 10))), min_size


def populate_bullets(tf, bullet_points, size, color, space_after=6):
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for idx, bullet in enumerate(bullet_points):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        style_paragraph(p, size=size, color=color)
        p.space_after = Pt(space_after)


def safe_add_textbox(slide, left_in, top_in, width_in, height_in, text="", auto_fit=True):
    left_in = max(0.08, left_in)
    top_in = max(0.08, top_in)
    width_in = min(width_in, SLIDE_W_IN - left_in - 0.1)
    height_in = min(height_in, SLIDE_H_IN - top_in - 0.08)
    return add_textbox(slide, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in), text=text, auto_fit=auto_fit)


def add_card(slide, left, top, width, height, fill_color, line_color=None, line_width=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(line_width)
    try:
        shape.adjustments[0] = 0.18
    except Exception:
        pass
    return shape


def add_footer(slide, theme_tokens, page_text: str):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(FOOTER_Y_IN), Inches(11.6), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = theme_tokens["line"]
    line.line.fill.background()
    _, tf = add_textbox(slide, Inches(11.7), Inches(6.9), Inches(0.75), Inches(0.22))
    p = tf.paragraphs[0]
    p.text = page_text
    style_paragraph(p, size=10, color=theme_tokens["muted"], align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, note_text: str):
    try:
        notes_text_frame = slide.notes_slide.notes_text_frame
        if notes_text_frame:
            notes_text_frame.text = (note_text or "").strip()[:1200]
    except Exception:
        pass


def render_cover_slide(prs: Presentation, ppt_data: dict, theme_tokens: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.24), theme_tokens["secondary"])
    add_full_rect(slide, Inches(0.78), Inches(1.15), Inches(0.2), Inches(4.7), theme_tokens["secondary"])
    title = clamp_text(ppt_data.get("title", "Selena Presentation").strip(), 32)
    title_size = adaptive_font_size(title, 30, 22, [(20, 26), (28, 23)])
    _, tf = safe_add_textbox(slide, 1.35, 1.28, 9.65, 1.1, auto_fit=True)
    p = tf.paragraphs[0]
    p.text = title
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])
    subtitle = clamp_text(ppt_data.get("subtitle") or "AI-generated presentation", 52)
    _, tf2 = safe_add_textbox(slide, 1.38, 2.38, 7.6, 0.75, auto_fit=True)
    p3 = tf2.paragraphs[0]
    p3.text = subtitle
    style_paragraph(p3, size=16, color=theme_tokens["muted"])
    _, tf4 = safe_add_textbox(slide, 1.38, 5.72, 4.3, 0.34, auto_fit=True)
    p4 = tf4.paragraphs[0]
    p4.text = clamp_text(ppt_data.get("author") or "Generated by Selena AI", 36)
    style_paragraph(p4, size=11, color=theme_tokens["muted"])
    return slide


def render_section_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.3), theme_tokens["primary"])
    add_card(slide, Inches(0.92), Inches(1.45), Inches(11.1), Inches(3.95), theme_tokens["accent"], theme_tokens["line"], 1.0)
    title = clamp_text(slide_data.get("title") or "Section", 28)
    _, tf = safe_add_textbox(slide, 1.32, 2.1, 9.0, 0.95, auto_fit=True)
    p = tf.paragraphs[0]
    p.text = title
    style_paragraph(p, size=28, bold=True, color=theme_tokens["primary"], align=PP_ALIGN.CENTER)
    summary = clamp_text(slide_data.get("summary") or "", 72)
    _, tf2 = safe_add_textbox(slide, 2.05, 3.2, 7.6, 0.8, auto_fit=True)
    p2 = tf2.paragraphs[0]
    p2.text = summary
    style_paragraph(p2, size=15, color=theme_tokens["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary)
    return slide


def render_bullet_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Untitled", 10.0, 0.62, 23, 18, max_len=34, bold=True)
    _, title_tf = safe_add_textbox(slide, 1.0, 0.72, 10.0, 0.62, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])
    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 9.95, 0.5, 15, 11, max_len=68)
    _, sub_tf = safe_add_textbox(slide, 1.0, 1.42, 9.95, 0.5, auto_fit=True)
    p2 = sub_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])
    add_card(slide, Inches(0.98), Inches(2.08), Inches(10.72), Inches(3.72), theme_tokens["accent"], theme_tokens["line"], 1.0)
    bullet_points, bullet_size, _ = fit_bullets_to_box(slide_data.get("bullets"), 9.1, 2.95, preferred_size=17, min_size=12, max_items=5, max_len=58, fallback_text=summary_text or "关键点")
    _, body_tf = safe_add_textbox(slide, 1.35, 2.42, 9.15, 3.05, auto_fit=True)
    populate_bullets(body_tf, bullet_points, bullet_size, theme_tokens["primary"], space_after=6)
    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide


def render_highlight_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    add_full_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.22), theme_tokens["secondary"])
    add_card(slide, Inches(0.95), Inches(1.08), Inches(11.2), Inches(4.9), theme_tokens["accent"], theme_tokens["line"], 1.2)
    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Key Message", 9.8, 0.55, 18, 13, max_len=26, bold=True)
    _, title_tf = safe_add_textbox(slide, 1.32, 1.42, 9.8, 0.55, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["secondary"])
    highlight_text, hl_size = fit_single_text_to_box(slide_data.get("highlight") or slide_data.get("summary") or "", 9.7, 1.6, 30, 18, max_len=62, bold=True)
    _, highlight_tf = safe_add_textbox(slide, 1.32, 2.1, 9.7, 1.6, auto_fit=True)
    hp = highlight_tf.paragraphs[0]
    hp.text = highlight_text
    style_paragraph(hp, size=hl_size, bold=True, color=theme_tokens["primary"], align=PP_ALIGN.CENTER)
    closing_text, close_size = fit_single_text_to_box(slide_data.get("closing") or slide_data.get("summary") or "", 8.8, 0.66, 14, 11, max_len=60)
    _, close_tf = safe_add_textbox(slide, 1.8, 4.45, 8.8, 0.66, auto_fit=True)
    cp = close_tf.paragraphs[0]
    cp.text = closing_text
    style_paragraph(cp, size=close_size, color=theme_tokens["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or highlight_text)
    return slide


def render_two_column_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Two Column", 10.3, 0.6, 22, 17, max_len=34, bold=True)
    _, title_tf = safe_add_textbox(slide, 0.95, 0.68, 10.3, 0.6, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])
    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.0, 0.5, 15, 11, max_len=62)
    _, sub_tf = safe_add_textbox(slide, 0.95, 1.38, 10.0, 0.5, auto_fit=True)
    p2 = sub_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])
    add_card(slide, Inches(0.95), Inches(2.08), Inches(5.2), Inches(3.95), theme_tokens["accent"], theme_tokens["line"], 1.0)
    add_card(slide, Inches(6.35), Inches(2.08), Inches(5.2), Inches(3.95), theme_tokens["accent"], theme_tokens["line"], 1.0)
    left_title, left_title_size = fit_single_text_to_box(slide_data.get("left_title") or "左侧观点", 4.3, 0.34, 16, 12, max_len=18, bold=True)
    left_points, left_size, _ = fit_bullets_to_box(slide_data.get("left_points") or (slide_data.get("bullets") or [])[:2], 4.3, 2.85, preferred_size=15, min_size=11, max_items=3, max_len=34, fallback_text=summary_text)
    _, left_tf = safe_add_textbox(slide, 1.15, 2.45, 4.35, 3.35, auto_fit=True)
    lp = left_tf.paragraphs[0]
    lp.text = left_title
    style_paragraph(lp, size=left_title_size, bold=True, color=theme_tokens["primary"])
    for point in left_points:
        pp = left_tf.add_paragraph()
        pp.text = f"• {point}"
        style_paragraph(pp, size=left_size, color=theme_tokens["primary"])
        pp.space_after = Pt(5)
    right_title, right_title_size = fit_single_text_to_box(slide_data.get("right_title") or "右侧解释", 4.3, 0.34, 16, 12, max_len=18, bold=True)
    right_points, right_size, _ = fit_bullets_to_box(slide_data.get("right_points") or (slide_data.get("bullets") or [])[2:], 4.3, 2.85, preferred_size=15, min_size=11, max_items=3, max_len=34, fallback_text=summary_text)
    _, right_tf = safe_add_textbox(slide, 6.5, 2.45, 4.35, 3.35, auto_fit=True)
    rp = right_tf.paragraphs[0]
    rp.text = right_title
    style_paragraph(rp, size=right_title_size, bold=True, color=theme_tokens["primary"])
    for point in right_points:
        pp = right_tf.add_paragraph()
        pp.text = f"• {point}"
        style_paragraph(pp, size=right_size, color=theme_tokens["primary"])
        pp.space_after = Pt(5)
    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide


def render_image_slide(prs, slide_data, theme_tokens, page_text, image_on_left=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Visual Slide", 10.6, 0.62, 22, 17, max_len=30, bold=True)
    _, title_tf = safe_add_textbox(slide, 0.82, 0.56, 10.6, 0.62, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])
    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.6, 0.5, 14, 11, max_len=60)
    _, summary_tf = safe_add_textbox(slide, 0.82, 1.22, 10.6, 0.5, auto_fit=True)
    p2 = summary_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])
    img_w = 5.0
    text_w = 4.9
    block_h = 3.85
    img_left = 0.82 if image_on_left else 7.51
    text_left = 6.05 if image_on_left else 0.82
    block_top = 1.95
    add_card(slide, Inches(img_left), Inches(block_top), Inches(img_w), Inches(block_h), theme_tokens["accent"], theme_tokens["line"], 1.0)
    image_path = slide_data.get("image_path") or ""
    if image_path and os.path.exists(image_path):
        add_picture_fit(slide, image_path, Inches(img_left), Inches(block_top), Inches(img_w), Inches(block_h), padding=Inches(0.10), mode="contain")
    else:
        _, placeholder_tf = safe_add_textbox(slide, img_left + 0.3, block_top + 1.45, img_w - 0.6, 0.45, auto_fit=True)
        pp = placeholder_tf.paragraphs[0]
        pp.text = clamp_text(slide_data.get("image_caption") or "视觉图像区域", 18)
        style_paragraph(pp, size=13, bold=True, color=theme_tokens["secondary"], align=PP_ALIGN.CENTER)
    add_card(slide, Inches(text_left), Inches(block_top), Inches(text_w), Inches(block_h), theme_tokens["accent"], theme_tokens["line"], 1.0)
    caption_text = (slide_data.get("image_caption") or "").strip()
    bullets = _dedupe_preserve_order(slide_data.get("bullets") or [])
    bullets = [b for b in bullets if _norm_text(b) not in {_norm_text(caption_text), _norm_text(summary_text)}]
    _, text_tf = safe_add_textbox(slide, text_left + 0.26, block_top + 0.22, text_w - 0.52, block_h - 0.42, auto_fit=True)
    if caption_text:
        cap_text, cap_size = fit_single_text_to_box(caption_text, text_w - 0.75, 0.54, 16, 12, max_len=28, bold=True)
        caption = text_tf.paragraphs[0]
        caption.text = cap_text
        style_paragraph(caption, size=cap_size, bold=True, color=theme_tokens["secondary"])
        caption.space_after = Pt(4)
    if summary_text:
        summary_fit, body_size = fit_single_text_to_box(summary_text, text_w - 0.75, 0.85, 13, 10, max_len=70)
        sp = text_tf.add_paragraph() if caption_text else text_tf.paragraphs[0]
        sp.text = summary_fit
        style_paragraph(sp, size=body_size, color=theme_tokens["muted"])
        sp.space_after = Pt(6)
    else:
        body_size = 12
    bullet_points, image_bullet_size, _ = fit_bullets_to_box(bullets, text_w - 0.75, 1.95, preferred_size=max(11, body_size), min_size=9, max_items=3, max_len=30)
    for point in bullet_points:
        pp = text_tf.add_paragraph()
        pp.text = f"• {point}"
        style_paragraph(pp, size=image_bullet_size, color=theme_tokens["primary"])
        pp.space_after = Pt(4)
    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text or caption_text)
    return slide


def render_compare_slide(prs, slide_data, theme_tokens, page_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, 0, 0, prs.slide_width, prs.slide_height, theme_tokens["light"])
    title_text, title_size = fit_single_text_to_box(slide_data.get("title") or "Compare", 10.4, 0.7, 23, 18, max_len=34, bold=True)
    _, title_tf = safe_add_textbox(slide, 0.95, 0.68, 10.4, 0.7, auto_fit=True)
    p = title_tf.paragraphs[0]
    p.text = title_text
    style_paragraph(p, size=title_size, bold=True, color=theme_tokens["primary"])
    summary_text, sum_size = fit_single_text_to_box(slide_data.get("summary") or "", 10.2, 0.56, 15, 11, max_len=60)
    _, sub_tf = safe_add_textbox(slide, 0.95, 1.42, 10.2, 0.56, auto_fit=True)
    p2 = sub_tf.paragraphs[0]
    p2.text = summary_text
    style_paragraph(p2, size=sum_size, color=theme_tokens["muted"])
    add_full_rect(slide, Inches(0.95), Inches(2.08), Inches(4.95), Inches(0.46), theme_tokens["primary"])
    add_full_rect(slide, Inches(6.28), Inches(2.08), Inches(4.95), Inches(0.46), theme_tokens["secondary"])
    _, l_head_tf = safe_add_textbox(slide, 1.12, 2.15, 4.55, 0.26, auto_fit=True)
    lp = l_head_tf.paragraphs[0]
    lp.text = clamp_text(slide_data.get("left_title") or "方案 A", 14)
    style_paragraph(lp, size=13, bold=True, color=theme_tokens["light"])
    _, r_head_tf = safe_add_textbox(slide, 6.45, 2.15, 4.55, 0.26, auto_fit=True)
    rp = r_head_tf.paragraphs[0]
    rp.text = clamp_text(slide_data.get("right_title") or "方案 B", 14)
    style_paragraph(rp, size=13, bold=True, color=theme_tokens["light"])
    add_card(slide, Inches(0.95), Inches(2.65), Inches(4.95), Inches(3.55), theme_tokens["accent"])
    add_card(slide, Inches(6.28), Inches(2.65), Inches(4.95), Inches(3.55), theme_tokens["accent"])
    left_points, left_size, _ = fit_bullets_to_box(slide_data.get("left_points"), 4.2, 2.9, preferred_size=14, min_size=10, max_items=4, max_len=30, fallback_text=summary_text)
    _, left_tf = safe_add_textbox(slide, 1.22, 2.95, 4.18, 2.95, auto_fit=True)
    populate_bullets(left_tf, left_points, left_size, theme_tokens["primary"], space_after=4)
    right_points, right_size, _ = fit_bullets_to_box(slide_data.get("right_points"), 4.2, 2.9, preferred_size=14, min_size=10, max_items=4, max_len=30, fallback_text=summary_text)
    _, right_tf = safe_add_textbox(slide, 6.55, 2.95, 4.18, 2.95, auto_fit=True)
    populate_bullets(right_tf, right_points, right_size, theme_tokens["primary"], space_after=4)
    add_footer(slide, theme_tokens, page_text)
    add_speaker_notes(slide, slide_data.get("speaker_note") or summary_text)
    return slide


def build_pptx_file(ppt_data: dict) -> BytesIO:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed.")
    ppt_data = normalize_ppt_data(ppt_data)
    ppt_data = ensure_slide_images(ppt_data, topic=ppt_data.get("title") or "Presentation", max_images=2)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    theme_tokens = get_theme_tokens(ppt_data.get("theme", "selena"))
    render_cover_slide(prs, ppt_data, theme_tokens)
    total_content_slides = len(ppt_data.get("slides", []))
    for idx, slide_data in enumerate(ppt_data.get("slides", []), start=1):
        slide_type = slide_data.get("type") or "bullet"
        page_text = f"{idx}/{total_content_slides}"
        if slide_type == "section":
            render_section_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "highlight":
            render_highlight_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "two_column":
            render_two_column_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "compare":
            render_compare_slide(prs, slide_data, theme_tokens, page_text)
        elif slide_type == "image_left":
            render_image_slide(prs, slide_data, theme_tokens, page_text, image_on_left=True)
        elif slide_type == "image_right":
            render_image_slide(prs, slide_data, theme_tokens, page_text, image_on_left=False)
        else:
            render_bullet_slide(prs, slide_data, theme_tokens, page_text)
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def summarize_ppt_for_chat(ppt_data: dict) -> str:
    slides = ppt_data.get("slides") or []
    overview = "\n".join([f"- {s.get('title', '未命名页面')}（{s.get('type', 'bullet')}）" for s in slides[:6]])
    image_count = int(ppt_data.get("generated_image_count") or 0)
    image_line = (
        f"**配图**：已自动生成并插入 {image_count} 张页面图片\n" if image_count
        else "**配图**：这次未插入图片，已按文字版式生成\n"
    )
    return (
        f"我已经按 **{ppt_data.get('theme', 'selena').title()}** 风格生成了一份可下载的 PPT。\n\n"
        f"**标题**：{ppt_data.get('title', '未命名演示')}\n"
        f"**副标题**：{ppt_data.get('subtitle', '')}\n"
        f"**页数**：{len(slides) + 1} 页（含封面）\n"
        f"{image_line}\n"
        f"**目录预览**：\n{overview}\n\n"
        f"这版会根据主题自动挑选 1 到 2 页做视觉化表达，并把生成的图片直接嵌入到 PPT 版式里。"
    )
